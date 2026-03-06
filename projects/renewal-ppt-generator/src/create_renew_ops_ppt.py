'''
Enhanced renewals PPT generator.

Generates Cisco renewal opportunity PowerPoints for product and service renewals
from CS Console Excel exports.

Usage:
    python create_renew_ops_ppt.py <initial_fy> <final_fy> <excel_filename> [--min-atr <value>] [--template-pptx <file>]

Example:
    python create_renew_ops_ppt.py Q1FY26 Q3FY26 .\\data\\renewals\\renewals.xlsx --min-atr 100

Inputs:
    initial_fy      Fiscal quarter start (QQFYXX), e.g. Q1FY26
    final_fy        Fiscal quarter end (QQFYXX), e.g. Q3FY26
    excel_filename  Renewals export file (.xlsx)
    --min-atr       Optional minimum aggregated Expected ATR ($000s)
    --template-pptx Optional PowerPoint template (.pptx)

Outputs:
    - <input>_product_<FY-range>[_MIN_ATR_###K].pptx
    - <input>_service_<FY-range>[_MIN_ATR_###K].pptx

Behavior:
    - Supports single-customer or all-customer input files
    - Aggregates rows by Deal Id for cleaner totals
    - Adds summary slides for All Customers and per-customer views
    - Adds account timelines and monthly timeline views
'''


import argparse
import os
import sys
from datetime import datetime
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
from matplotlib.dates import MonthLocator, DateFormatter
import matplotlib.patches as mpatches
import matplotlib.lines as mlines
import numpy as np
import io

# --- CONSTANTS ---
CISCO_FY_QUARTERS = {
    'Q1': (8, 1),   # August 1
    'Q2': (11, 1),  # November 1
    'Q3': (2, 1),   # February 1
    'Q4': (5, 1)    # May 1
}
REQUIRED_COLUMNS = [
    'Account ARR ($000s)', 'Account Name', 'CX Upsell/PMG', 'Close Date', 'Customer Name',
    'Customer Pulse', 'Deal Id', 'Deal Pulse', 'Expected ATR ($000s)', 'Expiration Date',
    'Expiration Quarter', 'Linked/Related', 'Linked/Related Deals', 'Opportunity Name',
    'Opportunity Owner', 'Opportunity Status', 'Prior ATR ($000s)', 'Product Amount (TCV) ($000s)',
    'Renewal Risk', 'Service Amount (TCV) ($000s)', 'Stage', 'Success Priority'
]
COLUMN_ALIASES = {
    'Renewal Risk': 'RenewalLine Risk'
}
BASE_SLIDE_COLUMNS = [
    'Account Name', 'Deal Id', 'Opportunity Name', 'CX Upsell/PMG',
    'Prior ATR ($000s)', 'Expected ATR ($000s)',
    'Product Amount (TCV) ($000s)',
    'Service Amount (TCV) ($000s)', 'Expiration Date',
    'Deal Pulse', 'Customer Pulse', 'Stage'  # Added Stage for monthly timelines
]
MAX_TABLE_ROWS_PER_SLIDE = 11
MAX_TIMELINES_PER_SLIDE = 16
MAX_MONTHLY_TIMELINES_PER_SLIDE = 16  # Maximum timeline rows per monthly slide
TIMELINE_SPACING = 0.4  # inches between timelines
BASE_HEIGHT = 1.5       # top/bottom padding for figure in inches
MIN_DAYS_SEPARATION = 3  # Minimum days between deals on the same timeline row

# Set your base circle size (e.g., current size used for $1 ATR)
MIN_CIRCLE_SIZE = 4
MAX_CIRCLE_SIZE = 16  # or MIN_CIRCLE_SIZE * 4

# Color name to RGB mapping for python-pptx and matplotlib
COLOR_RGB = {
    'black': RGBColor(0,0,0),
    'red': RGBColor(255,0,0),
    'green': RGBColor(0,128,0),
    'yellow': RGBColor(255,215,0),
    'blue': RGBColor(0,0,255),
    'orange': RGBColor(255,140,0),
    'purple': RGBColor(128,0,128),
    'grey': RGBColor(128,128,128),
}
DEFAULT_COLOR = 'black'

EMU_PER_INCH = 914400
SLIDE_MARGIN_IN = 0.35
CONTENT_BOTTOM_MARGIN_IN = 0.3
CONTENT_TITLE_HEIGHT_IN = 0.55
CONTENT_TITLE_GAP_IN = 0.12
TABLE_HEADER_HEIGHT_IN = 0.30
TABLE_ROW_HEIGHT_IN = 0.24


def remove_all_placeholders(slide):
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            sp = shape
            sp.element.getparent().remove(sp.element)


def get_slide_content_bounds(prs):
    left = Inches(SLIDE_MARGIN_IN)
    top = Inches(SLIDE_MARGIN_IN)
    right = Inches(SLIDE_MARGIN_IN)
    bottom = Inches(CONTENT_BOTTOM_MARGIN_IN)
    width = max(int(prs.slide_width - left - right), Inches(4))
    height = max(int(prs.slide_height - top - bottom), Inches(2))
    return left, top, width, height


def get_content_body_bounds(prs):
    left, top, width, content_height = get_slide_content_bounds(prs)
    title_height = Inches(CONTENT_TITLE_HEIGHT_IN)
    title_gap = Inches(CONTENT_TITLE_GAP_IN)
    body_top = top + title_height + title_gap
    body_height = max(int(content_height - title_height - title_gap), Inches(1.6))
    return left, top, width, title_height, body_top, body_height


def add_content_title(slide, prs, text, font_size=16):
    left, top, width, title_height, body_top, body_height = get_content_body_bounds(prs)
    textbox = slide.shapes.add_textbox(left, top, width, title_height)
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.text = text
    tf.paragraphs[0].font.size = Pt(font_size)
    tf.paragraphs[0].alignment = 1
    return left, body_top, width, body_height


def get_table_rows_per_slide(content_height_emu, include_totals=False):
    content_height_in = content_height_emu / EMU_PER_INCH
    reserved_in = TABLE_HEADER_HEIGHT_IN + (TABLE_ROW_HEIGHT_IN if include_totals else 0.0)
    usable_in = max(content_height_in - reserved_in, TABLE_ROW_HEIGHT_IN)
    return max(1, int(usable_in / TABLE_ROW_HEIGHT_IN))


def get_table_font_sizes(col_count):
    if col_count >= 10:
        return Pt(7), Pt(6)
    if col_count >= 8:
        return Pt(8), Pt(6)
    return Pt(9), Pt(7)


def apply_table_dimensions(table, total_width_emu, total_height_emu, first_col_ratio=0.30):
    col_count = len(table.columns)
    row_count = len(table.rows)
    if col_count > 1:
        first_col_width = int(total_width_emu * first_col_ratio)
        remaining = max(total_width_emu - first_col_width, 0)
        other_col_width = int(remaining / (col_count - 1)) if col_count > 1 else remaining
        table.columns[0].width = first_col_width
        for idx in range(1, col_count):
            table.columns[idx].width = other_col_width
    row_height = max(int(total_height_emu / max(row_count, 1)), Inches(0.18))
    for row in table.rows:
        row.height = row_height


def add_fitted_picture(slide, img_stream, left, top, max_width, max_height, image_aspect_ratio):
    if max_width <= 0 or max_height <= 0:
        return
    image_aspect = image_aspect_ratio if image_aspect_ratio and image_aspect_ratio > 0 else (16.0 / 9.0)
    box_aspect = max_width / max_height
    if image_aspect >= box_aspect:
        draw_width = max_width
        draw_height = int(draw_width / image_aspect)
    else:
        draw_height = max_height
        draw_width = int(draw_height * image_aspect)
    draw_left = left + int((max_width - draw_width) / 2)
    draw_top = top + int((max_height - draw_height) / 2)
    slide.shapes.add_picture(img_stream, draw_left, draw_top, width=draw_width, height=draw_height)

def format_currency(value):
    """Format currency values as integers without .0"""
    try:
        if pd.isnull(value):
            return ""
        val = float(value)
        if val == 0:
            return ""
        return f"${int(round(val)):,}"
    except (ValueError, TypeError):
        return ""

def safe_int_format(value):
    """Safely convert value to integer format, handling strings and nulls"""
    try:
        if pd.isnull(value):
            return ""
        val = float(value)
        if val == 0:
            return ""
        return f"{int(round(val)):,}"
    except (ValueError, TypeError):
        return str(value) if value else ""

def robust_validate_fy_quarter(fy_str):
    if not isinstance(fy_str, str):
        print(f"Error: Fiscal quarter must be a string, got {type(fy_str)}", file=sys.stderr)
        return False
    if len(fy_str) != 6:
        print(f"Error: Fiscal quarter '{fy_str}' must have 6 characters (e.g., Q1FY26).", file=sys.stderr)
        return False
    quarter = fy_str[:2]
    if quarter not in CISCO_FY_QUARTERS:
        print(f"Error: Quarter '{quarter}' in '{fy_str}' is invalid. Expected one of {list(CISCO_FY_QUARTERS)}.", file=sys.stderr)
        return False
    if fy_str[2:4] != 'FY':
        print(f"Error: Fiscal quarter '{fy_str}' must contain 'FY' after the quarter.", file=sys.stderr)
        return False
    if not fy_str[4:].isdigit():
        print(f"Error: Year '{fy_str[4:]}' in '{fy_str}' is not numeric.", file=sys.stderr)
        return False
    return True

def robust_check_excel_file(filename):
    if not isinstance(filename, str):
        print("Error: Excel filename must be a string.", file=sys.stderr)
        return False
    if not filename.lower().endswith('.xlsx'):
        print("Error: Excel filename must end with '.xlsx'.", file=sys.stderr)
        return False
    if not os.path.isfile(filename):
        print(f"Error: File '{filename}' does not exist.", file=sys.stderr)
        return False
    return True

def robust_check_template_file(filename):
    if filename is None:
        return True
    if not isinstance(filename, str):
        print("Error: Template filename must be a string.", file=sys.stderr)
        return False
    if not filename.lower().endswith('.pptx'):
        print("Error: Template file must end with '.pptx'. If you have a .potx, save/convert it to .pptx first.", file=sys.stderr)
        return False
    if not os.path.isfile(filename):
        print(f"Error: Template file '{filename}' does not exist.", file=sys.stderr)
        return False
    return True

def validate_fy_quarter(fy_str):
    quarter = fy_str[:2]
    year = int('20' + fy_str[-2:])
    month, day = CISCO_FY_QUARTERS[quarter]
    if quarter == 'Q1':
        start = datetime(year - 1, month, day)
        end = datetime(year - 1, 10, 31)
    elif quarter == 'Q2':
        start = datetime(year - 1, month, day)
        end = datetime(year, 1, 31)
    elif quarter == 'Q3':
        start = datetime(year, month, day)
        end = datetime(year, 4, 30)
    elif quarter == 'Q4':
        start = datetime(year, month, day)
        end = datetime(year, 7, 31)
    else:
        raise ValueError(f"Unknown quarter: {quarter}")
    return start, end

def get_fy_range(start_fy, end_fy):
    start_dt, _ = validate_fy_quarter(start_fy)
    _, end_dt = validate_fy_quarter(end_fy)
    return start_dt, end_dt

def get_file_creation_date(path):
    t = os.path.getctime(path)
    return datetime.fromtimestamp(t).strftime('%Y-%m-%d')

def check_excel_file(filename):
    if not os.path.isfile(filename):
        raise FileNotFoundError(f"Input Excel file '{filename}' does not exist.")
    try:
        df = pd.read_excel(filename)
    except Exception as e:
        raise ValueError(f"Could not read Excel file: {e}")
    for expected, actual in COLUMN_ALIASES.items():
        if expected not in df.columns and actual in df.columns:
            print(f"Info: Using '{actual}' as '{expected}'")
            df[expected] = df[actual]
    missing = [col for col in REQUIRED_COLUMNS if col not in df.columns]
    if missing:
        raise ValueError(f"Excel file missing required columns: {missing}")
    return df

def get_pulse_color(pulse_val):
    """
    Extract color name from Deal Pulse or Customer Pulse.
    Accepts values like '1 - Yellow', '2- yellow ', 'NA'.
    Returns color name as lower case, or 'black' if not found.
    """
    if isinstance(pulse_val, str) and pulse_val.strip().upper() != 'NA':
        # Find the part after the dash
        parts = pulse_val.split('-')
        if len(parts) == 2:
            color_name = parts[1].strip().lower()
            return color_name if color_name in COLOR_RGB else DEFAULT_COLOR
    return DEFAULT_COLOR


def get_customer_names(df):
    if 'Customer Name' not in df.columns:
        return []
    customer_names = (
        df['Customer Name']
        .dropna()
        .astype(str)
        .str.strip()
    )
    customer_names = customer_names[customer_names != '']
    return sorted(customer_names.unique().tolist())


def format_customer_scope(customer_names, max_names=4):
    if not customer_names:
        return 'Selected Customers'
    if len(customer_names) <= max_names:
        return ', '.join(customer_names)
    return f"{', '.join(customer_names[:max_names])} (+{len(customer_names) - max_names} more)"


def add_summary_table_slide(prs, filtered, fy_start, fy_end, value_col, title_prefix, summary_scope='All Customers'):
    if filtered.empty:
        return None
    months = pd.date_range(start=fy_start, end=fy_end, freq='MS')
    month_labels = [m.strftime('%b %Y') for m in months]
    account_names = sorted(filtered['Account Name'].unique())
    filtered = filtered.copy()
    filtered['Month'] = filtered['Expiration Date'].dt.to_period('M')
    table_data = pd.DataFrame(0.0, index=account_names, columns=month_labels)
    for account in account_names:
        acc_df = filtered[filtered['Account Name'] == account]
        for m, label in zip(months, month_labels):
            month_val = acc_df[(acc_df['Expiration Date'].dt.month == m.month) & (acc_df['Expiration Date'].dt.year == m.year)]
            total = month_val[value_col].sum()
            table_data.loc[account, label] = total

    # Add Total column and row
    table_data['Total ($000s)'] = table_data.sum(axis=1)
    total_row = table_data.sum(axis=0)
    table_cols = list(month_labels) + ['Total ($000s)']

    _left, _top, _width, _title_height, _body_top, body_height = get_content_body_bounds(prs)
    rows_per_slide = min(MAX_TABLE_ROWS_PER_SLIDE, get_table_rows_per_slide(body_height, include_totals=True))
    header_font, body_font = get_table_font_sizes(len(table_cols) + 1)
    total_accounts = list(table_data.index)
    num_pages = (len(total_accounts) + rows_per_slide - 1) // rows_per_slide

    last_slide = None
    for page_idx, start_row in enumerate(range(0, len(total_accounts), rows_per_slide), start=1):
        page_accounts = total_accounts[start_row:start_row + rows_per_slide]
        is_last_page = start_row + rows_per_slide >= len(total_accounts)
        title_suffix = f" (page {page_idx})" if num_pages > 1 else ""

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        remove_all_placeholders(slide)
        table_left, table_top, table_width, table_height = add_content_title(
            slide,
            prs,
            f"{title_prefix} Summary - {summary_scope} - Expected ATR Aggregated by Month ($000s){title_suffix}",
            font_size=16,
        )

        total_rows = 1 + len(page_accounts) + (1 if is_last_page else 0)
        table = slide.shapes.add_table(total_rows, len(table_cols) + 1, table_left, table_top, table_width, table_height).table
        apply_table_dimensions(table, table_width, table_height, first_col_ratio=0.28)

        table.cell(0, 0).text = "Account Name"
        table.cell(0, 0).text_frame.paragraphs[0].font.size = header_font
        table.cell(0, 0).text_frame.paragraphs[0].font.bold = True
        for j, label in enumerate(table_cols):
            cell = table.cell(0, j + 1)
            cell.text = label
            cell.text_frame.paragraphs[0].font.size = header_font
            cell.text_frame.paragraphs[0].font.bold = True

        for i, account in enumerate(page_accounts, start=1):
            table.cell(i, 0).text = account
            table.cell(i, 0).text_frame.paragraphs[0].font.size = body_font
            table.cell(i, 0).text_frame.paragraphs[0].font.bold = True
            for j, label in enumerate(table_cols):
                val = table_data.loc[account, label]
                cell = table.cell(i, j + 1)
                cell.text = format_currency(val)
                cell.text_frame.paragraphs[0].font.size = body_font
                if label == 'Total ($000s)':
                    cell.text_frame.paragraphs[0].font.bold = True

        if is_last_page:
            total_idx = len(page_accounts) + 1
            table.cell(total_idx, 0).text = "Total ($000s)"
            table.cell(total_idx, 0).text_frame.paragraphs[0].font.size = body_font
            table.cell(total_idx, 0).text_frame.paragraphs[0].font.bold = True
            for j, label in enumerate(table_cols):
                val = total_row[label]
                cell = table.cell(total_idx, j + 1)
                cell.text = format_currency(val)
                cell.text_frame.paragraphs[0].font.size = body_font
                cell.text_frame.paragraphs[0].font.bold = True
        last_slide = slide
    return last_slide

def add_table_slides(prs, filtered, columns, title_prefix, account_col='Account Name'):
    for account in filtered[account_col].unique():
        acc_table_df = filtered[filtered[account_col] == account][columns].copy()
                
        # Format currency columns to remove .0
        for col in ['Prior ATR ($000s)', 'Expected ATR ($000s)', 'Product Amount (TCV) ($000s)', 'Service Amount (TCV) ($000s)']:
            if col in acc_table_df.columns:
                acc_table_df[col] = acc_table_df[col].apply(safe_int_format)
        
        num_rows = acc_table_df.shape[0]
        _left, _top, _width, _title_height, _body_top, body_height = get_content_body_bounds(prs)
        rows_per_slide = min(MAX_TABLE_ROWS_PER_SLIDE, get_table_rows_per_slide(body_height, include_totals=False))
        num_slides = (num_rows // rows_per_slide) + (1 if num_rows % rows_per_slide else 0)
        for slide_idx, start_row in enumerate(range(0, num_rows, rows_per_slide)):
            chunk = acc_table_df.iloc[start_row:start_row+rows_per_slide]
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            remove_all_placeholders(slide)
            table_left, table_top, table_width, table_height = add_content_title(
                slide,
                prs,
                f"{account} {title_prefix} Opportunities" + (f" (page {slide_idx+1})" if num_slides > 1 else ""),
                font_size=15,
            )
            rows, cols = chunk.shape
            table = slide.shapes.add_table(rows+1, cols, table_left, table_top, table_width, table_height).table
            apply_table_dimensions(table, table_width, table_height, first_col_ratio=0.24)
            header_font, body_font = get_table_font_sizes(cols)
            for j, col in enumerate(chunk.columns):
                cell = table.cell(0, j)
                cell.text = str(col)
                cell.text_frame.paragraphs[0].font.size = header_font
                cell.text_frame.paragraphs[0].font.bold = True
            for i, row in enumerate(chunk.values):
                for j, colname in enumerate(chunk.columns):
                    cell = table.cell(i+1, j)
                    val = "" if pd.isnull(row[j]) else row[j].strftime('%Y-%m-%d') if isinstance(row[j], pd.Timestamp) else str(row[j])
                    cell.text = val
                    cell.text_frame.paragraphs[0].font.size = body_font



def add_timeline_legend(ax):
    # Example pulse colors, adjust to match your actual color mapping if needed
    pulse_colors = {
        'green': '#34A853',    # Example RGB value, use your own
        'yellow': '#FBBC05',
        'red': '#EA4335',
        'grey': '#B0B0B0'      # Standard Grey
    }
    # Use Line2D for circles in the legend
    legend_elements = [
        mlines.Line2D([], [], color=pulse_colors['green'], marker='o', linestyle='None', markersize=9, label='Deal Pulse color & ATR size in circle'),
        mpatches.Patch(facecolor='white', edgecolor='gray', label='Customer Pulse color in Deal Id text', alpha=0.0)
    ]
    # Place legend below the plot, centered
    ax.legend(
        handles=legend_elements,
        loc='lower center',
        bbox_to_anchor=(0.5, 1.02),
        ncol=2,
        fontsize=9,
        frameon=True
    )


def get_circle_size(atr, min_atr, max_atr, min_circle_size=MIN_CIRCLE_SIZE, max_circle_size=MAX_CIRCLE_SIZE):
        if max_atr == min_atr:
            return min_circle_size
        norm = (atr - min_atr) / (max_atr - min_atr)
        return min_circle_size + norm * (max_circle_size - min_circle_size)

def get_atr_min_max(df, account_name, atr_column='Expected ATR ($000s)'):
    """
    Returns min_atr (minimum, at least 1) and max_atr for all deals of a given Account Name.
    df: pandas DataFrame with all deals.
    account_name: string, the Account Name to filter by.
    atr_column: the column name for Expected ATR ($000s).
    """
    # Filter for the specified Account Name
    deals = df[df['Account Name'] == account_name]
    # Drop missing or non-numeric values
    atr_values = pd.to_numeric(deals[atr_column], errors='coerce').dropna()
    if atr_values.empty:
        # Return None or sensible defaults if no valid deals
        return 1.0, 1.0
    min_atr = max(atr_values.min(), 1)
    max_atr = atr_values.max()
    return min_atr, max_atr

def add_timeline_slides(prs, filtered, fy_start, fy_end, title_prefix, total_atr_by_deal, account_col='Account Name'):
    for account in filtered[account_col].unique():
        acc_df = filtered[filtered[account_col] == account].copy()
        acc_df = acc_df.sort_values('Expiration Date')

        min_atr = acc_df['Expected ATR ($000s)'].min()
        max_atr = acc_df['Expected ATR ($000s)'].max()

        timelines = []
        timeline_assignments = []
        for idx, row in acc_df.iterrows():
            placed = False
            for t_idx, timeline in enumerate(timelines):
                if (row['Expiration Date'] - timeline[-1]['Expiration Date']).days >= 15:
                    timeline.append(row)
                    timeline_assignments.append(t_idx)
                    placed = True
                    break
            if not placed:
                timelines.append([row])
                timeline_assignments.append(len(timelines)-1)
        
        acc_df['timeline_row'] = timeline_assignments
        n_timelines = len(timelines)
        
        for slide_idx, start_timeline in enumerate(range(0, n_timelines, MAX_TIMELINES_PER_SLIDE)):
            end_timeline = start_timeline + MAX_TIMELINES_PER_SLIDE
            
            slide_plot_df = acc_df[
                (acc_df['timeline_row'] >= start_timeline) & 
                (acc_df['timeline_row'] < end_timeline)
            ].copy()

            if slide_plot_df.empty:
                continue

            min_row = slide_plot_df['timeline_row'].min()
            slide_plot_df['plot_row'] = slide_plot_df['timeline_row'] - min_row
            n_rows_this_slide = slide_plot_df['plot_row'].max() + 1
            
            fig_height = BASE_HEIGHT + TIMELINE_SPACING * n_rows_this_slide
            fig_width = 10
            fig, ax = plt.subplots(figsize=(fig_width, fig_height))
            
            for i in range(n_rows_this_slide):
                ax.hlines(y=i, xmin=fy_start, xmax=fy_end, color='tab:blue', linewidth=2)
            
            for _, row in slide_plot_df.iterrows():
                x = row['Expiration Date']
                y = row['plot_row']
                deal_id = str(row['Deal Id'])
                stage = str(row.get('Stage', '')).strip()
                
                deal_pulse_color = get_pulse_color(row.get('Deal Pulse', 'NA')).lower()
                customer_pulse_color = get_pulse_color(row.get('Customer Pulse', 'NA')).lower()
                
                atr = float(row['Expected ATR ($000s)'])
                circle_size = get_circle_size(atr, min_atr, max_atr)
                
                ax.plot(x, y, 'o', color=deal_pulse_color, markersize=circle_size)
                
                is_closed_won = stage.lower() == '6 - closed won'
                
                if is_closed_won:
                    ax.text(
                        x, y + 0.12, deal_id,
                        fontsize=8, ha='center', va='bottom', rotation=20, fontweight='bold',
                        color='black',
                        bbox=dict(facecolor='lightgreen', edgecolor='green', linewidth=1.5, boxstyle='round,pad=0.3')
                    )
                elif customer_pulse_color == 'yellow':
                    ax.text(
                        x, y + 0.12, deal_id,
                        fontsize=8, ha='center', va='bottom', rotation=20, fontweight='bold',
                        color='black',
                        bbox=dict(facecolor='yellow', edgecolor='none', boxstyle='round,pad=0.1')
                    )
                else:
                    mpl_text_color = customer_pulse_color if customer_pulse_color in COLOR_RGB else 'black'
                    if mpl_text_color == 'yellow':
                        mpl_text_color = 'black'
                    ax.text(
                        x, y + 0.12, deal_id,
                        fontsize=8, ha='center', va='bottom', rotation=20, fontweight='bold',
                        color=mpl_text_color
                    )            
            
            months = pd.date_range(start=fy_start, end=fy_end, freq='MS')
            for month_start in months:
                ax.vlines(month_start, ymin=-0.5, ymax=n_rows_this_slide-0.5, color='lightgrey', linestyle='--', linewidth=1)
            
            ax.set_yticks([])
            ax.set_xlim([fy_start, fy_end])
            ax.set_ylim(-0.5, n_rows_this_slide - 0.5)
            ax.set_title("")
            ax.set_xlabel('Expiration Date')
            ax.xaxis.set_major_locator(MonthLocator())
            ax.xaxis.set_major_formatter(DateFormatter('%b\n%Y'))

            add_timeline_legend(ax)

            plt.tight_layout()
            img_stream = io.BytesIO()
            plt.savefig(img_stream, format='png')
            plt.close(fig)
            img_stream.seek(0)
            image_aspect_ratio = fig_width / fig_height if fig_height > 0 else (16.0 / 9.0)
            
            timeline_slide = prs.slides.add_slide(prs.slide_layouts[1])
            remove_all_placeholders(timeline_slide)
            pic_left, pic_top, pic_width, pic_height = add_content_title(
                timeline_slide,
                prs,
                f"{account} {title_prefix} Opportunities Timeline" + (f" (page {slide_idx+1})" if n_timelines > MAX_TIMELINES_PER_SLIDE else ""),
                font_size=15,
            )
            add_fitted_picture(timeline_slide, img_stream, pic_left, pic_top, pic_width, pic_height, image_aspect_ratio)

            # --- Add deal details to Notes with TOTAL ATR ---
            notes_lines = ["Detailed Entries:"]
            
            slide_plot_df['Month'] = slide_plot_df['Expiration Date'].dt.strftime('%B %Y')
            for month in sorted(slide_plot_df['Month'].unique(), key=lambda x: datetime.strptime(x, '%B %Y')):
                notes_lines.append(f"{month}:")
                month_deals = slide_plot_df[slide_plot_df['Month'] == month]
                
                for _, deal in month_deals.iterrows():
                    deal_id = deal['Deal Id']
                    deal_id_str = str(deal_id)
                    account_name = deal.get('Account Name', "")
                    # This PPT's ATR (product or service portion)
                    atr_this_ppt = int(round(float(deal['Expected ATR ($000s)'])))
                    # Total ATR across both product and service
                    atr_total = int(round(total_atr_by_deal.get(deal_id, atr_this_ppt)))
                    opp_name = deal.get('Opportunity Name', "")
                    stage = deal.get('Stage', "")
                    upsell = deal.get('CX Upsell/PMG', "")
                    detail = f"- {deal_id_str}: {account_name} | {atr_this_ppt} | {atr_total} | {opp_name} | {stage} | {upsell}"
                    notes_lines.append(detail)

            notes_text = '\n'.join(notes_lines)
            timeline_slide.notes_slide.notes_text_frame.text = notes_text

def add_monthly_timeline_slides(prs, filtered, fy_start, fy_end, title_prefix, total_atr_by_deal):
    """
    Create monthly timeline slides showing all accounts' renewal opportunities for each month.
    """
    agg_df = filtered.copy()
    agg_df['Month'] = agg_df['Expiration Date'].dt.to_period('M')
    
    months = pd.date_range(start=fy_start, end=fy_end, freq='MS')
    
    for month in months:
        month_period = month.to_period('M')
        month_df = agg_df[agg_df['Month'] == month_period].copy()
        
        if month_df.empty:
            continue
        
        min_atr = max(month_df['Expected ATR ($000s)'].min(), 1)
        max_atr = month_df['Expected ATR ($000s)'].max()
        
        month_df = month_df.sort_values('Expiration Date').reset_index(drop=True)
        
        timelines = []
        
        for idx, row in month_df.iterrows():
            placed = False
            current_date = row['Expiration Date']
            
            for timeline in timelines:
                can_fit = True
                for existing_deal in timeline:
                    days_diff = abs((current_date - existing_deal['Expiration Date']).days)
                    if days_diff < MIN_DAYS_SEPARATION:
                        can_fit = False
                        break
                
                if can_fit:
                    timeline.append(row.to_dict())
                    placed = True
                    break
            
            if not placed:
                timelines.append([row.to_dict()])
        
        n_rows = len(timelines)
        
        month_start = pd.Timestamp(month.year, month.month, 1)
        if month.month == 12:
            month_end = pd.Timestamp(month.year + 1, 1, 1) - pd.Timedelta(days=1)
        else:
            month_end = pd.Timestamp(month.year, month.month + 1, 1) - pd.Timedelta(days=1)
        
        num_slides = (n_rows // MAX_MONTHLY_TIMELINES_PER_SLIDE) + (1 if n_rows % MAX_MONTHLY_TIMELINES_PER_SLIDE else 0)
        
        for slide_idx in range(num_slides):
            start_row = slide_idx * MAX_MONTHLY_TIMELINES_PER_SLIDE
            end_row = min(start_row + MAX_MONTHLY_TIMELINES_PER_SLIDE, n_rows)
            slide_timelines = timelines[start_row:end_row]
            n_rows_this_slide = len(slide_timelines)
            
            fig_height = BASE_HEIGHT + TIMELINE_SPACING * max(n_rows_this_slide, 5)
            fig_width = 10
            fig, ax = plt.subplots(figsize=(fig_width, fig_height))
            
            for row_idx, timeline in enumerate(slide_timelines):
                y = row_idx
                
                ax.hlines(y=y, xmin=month_start, xmax=month_end, color='tab:blue', linewidth=2)
                
                for deal in timeline:
                    x = deal['Expiration Date']
                    deal_id = str(deal['Deal Id'])
                    stage = str(deal.get('Stage', '')).strip()
                    
                    deal_pulse_color = get_pulse_color(deal.get('Deal Pulse', 'NA')).lower()
                    customer_pulse_color = get_pulse_color(deal.get('Customer Pulse', 'NA')).lower()
                    
                    atr = float(deal['Expected ATR ($000s)'])
                    circle_size = get_circle_size(atr, min_atr, max_atr)
                    
                    ax.plot(x, y, 'o', color=deal_pulse_color, markersize=circle_size)
                    
                    is_closed_won = stage.lower() == '6 - closed won'
                    
                    if is_closed_won:
                        ax.text(
                            x, y + 0.12, deal_id,
                            fontsize=8, ha='center', va='bottom', rotation=20, fontweight='bold',
                            color='black',
                            bbox=dict(facecolor='lightgreen', edgecolor='green', linewidth=1.5, boxstyle='round,pad=0.3')
                        )
                    elif customer_pulse_color == 'yellow':
                        ax.text(
                            x, y + 0.12, deal_id,
                            fontsize=8, ha='center', va='bottom', rotation=20, fontweight='bold',
                            color='black',
                            bbox=dict(facecolor='yellow', edgecolor='none', boxstyle='round,pad=0.1')
                        )
                    else:
                        mpl_text_color = customer_pulse_color if customer_pulse_color in COLOR_RGB else 'black'
                        if mpl_text_color == 'yellow':
                            mpl_text_color = 'black'
                        ax.text(
                            x, y + 0.12, deal_id,
                            fontsize=8, ha='center', va='bottom', rotation=20, fontweight='bold',
                            color=mpl_text_color
                        )
            
            ax.set_yticks([])
            ax.set_xlim([month_start, month_end])
            ax.set_ylim(-0.5, n_rows_this_slide - 0.5)
            ax.set_title("")
            ax.set_xlabel('Expiration Date')
            ax.xaxis.set_major_locator(MonthLocator())
            ax.xaxis.set_major_formatter(DateFormatter('%b\n%Y'))
            
            add_timeline_legend(ax)
            
            plt.tight_layout()
            img_stream = io.BytesIO()
            plt.savefig(img_stream, format='png')
            plt.close(fig)
            img_stream.seek(0)
            image_aspect_ratio = fig_width / fig_height if fig_height > 0 else (16.0 / 9.0)
            
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            remove_all_placeholders(slide)
            page_suffix = f" (page {slide_idx + 1})" if num_slides > 1 else ""
            pic_left, pic_top, pic_width, pic_height = add_content_title(
                slide,
                prs,
                f"New Renewal Opportunities in {month.strftime('%B %Y')}{page_suffix}",
                font_size=16,
            )
            add_fitted_picture(slide, img_stream, pic_left, pic_top, pic_width, pic_height, image_aspect_ratio)
            
            # Add notes with TOTAL ATR
            notes_lines = ["Detailed Entries:", f"{month.strftime('%B %Y')}:"]
            
            slide_deals = []
            for timeline in slide_timelines:
                slide_deals.extend(timeline)
            
            slide_deals_df = pd.DataFrame(slide_deals).sort_values('Expiration Date')
            
            for _, row in slide_deals_df.iterrows():
                deal_id = row['Deal Id']
                deal_id_str = str(deal_id)
                account = row['Account Name']
                # This PPT's ATR
                atr_this_ppt = int(round(float(row['Expected ATR ($000s)'])))
                # Total ATR across both product and service
                atr_total = int(round(total_atr_by_deal.get(deal_id, atr_this_ppt)))
                opp_name = row.get('Opportunity Name', "")
                stage = row.get('Stage', "")
                upsell = row.get('CX Upsell/PMG', "")
                
                detail = f"- {deal_id_str}: {account} | {atr_this_ppt} | {atr_total} | {opp_name} | {stage} | {upsell}"
                notes_lines.append(detail)
            
            notes_text = '\n'.join(notes_lines)
            slide.notes_slide.notes_text_frame.text = notes_text

            
def create_ppt_for_filter(filtered, fy_start, fy_end, excel_filename, suffix, initial_fy, final_fy, file_creation_date, title_prefix, columns, min_atr, total_atr_by_deal, template_pptx=None):
    # PRE-AGGREGATE DATA BY DEAL ID BEFORE ANY SLIDE CREATION
    aggregated_by_deal = (
        filtered
        .groupby('Deal Id')
        .agg({
            'Customer Name': 'first',
            'Account Name': 'first',
            'Opportunity Name': 'first',
            'CX Upsell/PMG': 'first',
            'Prior ATR ($000s)': 'sum',
            'Expected ATR ($000s)': 'sum',
            'Product Amount (TCV) ($000s)': 'sum',
            'Service Amount (TCV) ($000s)': 'sum',
            'Expiration Date': 'first',
            'Deal Pulse': 'first',
            'Customer Pulse': 'first',
            'Stage': 'first'
        })
        .reset_index()
    )
    aggregated_by_deal['Customer Name'] = aggregated_by_deal['Customer Name'].fillna('').astype(str).str.strip()
    customer_names = get_customer_names(aggregated_by_deal)
    customer_scope = format_customer_scope(customer_names, max_names=5)
    
    prs = Presentation(template_pptx) if template_pptx else Presentation()
    # --- TITLE SLIDE ---
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = title_slide.shapes.title
    title_text = f"{title_prefix} Opportunities in {customer_scope}"
    if title_shape is not None:
        title_shape.text = title_text
    else:
        title_box = title_slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8.8), Inches(0.8))
        title_box.text_frame.text = title_text
    subtitle_text = (
        f"Date Range: {initial_fy} – {final_fy}\n"
        f"Fiscal Dates: {fy_start.strftime('%b %d, %Y')} – {fy_end.strftime('%b %d, %Y')}\n"
        f"File Created: {file_creation_date}"
    )
    if customer_names:
        subtitle_text += f"\nCustomers ({len(customer_names)}): {format_customer_scope(customer_names, max_names=8)}"
    else:
        subtitle_text += "\nCustomers: Not provided"
    if min_atr > 0:
        subtitle_text += f"\nMinimum Expected ATR: ${int(min_atr):,}"
    try:
        subtitle_shape = title_slide.placeholders[1]
        subtitle_shape.text = subtitle_text
        subtitle_tf = subtitle_shape.text_frame
    except Exception:
        subtitle_box = title_slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.5), Inches(1.8))
        subtitle_tf = subtitle_box.text_frame
        subtitle_tf.text = subtitle_text
    
    # Reduce subtitle font size
    for paragraph in subtitle_tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(14)
    
    # Add legend text box on title slide
    if min_atr > 0:
        # Position the legend at the bottom of the slide
        legend_left = Inches(1)
        legend_top = Inches(6)
        legend_width = Inches(8)
        legend_height = Inches(1)
        
        legend_box = title_slide.shapes.add_textbox(legend_left, legend_top, legend_width, legend_height)
        legend_frame = legend_box.text_frame
        legend_frame.word_wrap = True
        
        # Add legend title
        p = legend_frame.paragraphs[0]
        p.text = "Note:"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        
        # Add explanation
        p = legend_frame.add_paragraph()
        p.text = f"This presentation includes only deals with total aggregated Expected ATR >= ${int(min_atr)}K"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_before = Pt(6)
        
        # Add note about the two ATR values in notes
        p = legend_frame.add_paragraph()
        p.text = "Notes pages show: Deal portion ATR | Total Deal ATR (Product + Service)"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_before = Pt(3)
    
    # --- SUMMARY TABLE SLIDES ---
    value_col = 'Expected ATR ($000s)'
    add_summary_table_slide(
        prs, aggregated_by_deal, fy_start, fy_end, value_col, title_prefix, summary_scope='All Customers'
    )
    for customer in customer_names:
        customer_df = aggregated_by_deal[aggregated_by_deal['Customer Name'] == customer]
        add_summary_table_slide(
            prs, customer_df, fy_start, fy_end, value_col, title_prefix, summary_scope=customer
        )
    
    # --- Table and Timeline slides (pass total_atr_by_deal) ---
    add_table_slides(prs, aggregated_by_deal, columns, title_prefix)
    add_timeline_slides(prs, aggregated_by_deal, fy_start, fy_end, title_prefix, total_atr_by_deal)
    
    # --- Monthly Timeline slides (pass total_atr_by_deal) ---
    add_monthly_timeline_slides(prs, aggregated_by_deal, fy_start, fy_end, title_prefix, total_atr_by_deal)
    
    pptx_filename = os.path.splitext(excel_filename)[0] + suffix + '.pptx'
    prs.save(pptx_filename)
    print(f"Presentation saved as: {pptx_filename}")



def create_renewal_ppt(initial_fy, final_fy, excel_filename, min_atr=0, template_pptx=None):
    fy_start, fy_end = get_fy_range(initial_fy, final_fy)
    df = check_excel_file(excel_filename)
    if len(df) == 0:
        print("Warning: The input Excel file contains no data.")
        return
        
    df['Expiration Date'] = pd.to_datetime(df['Expiration Date'], errors='coerce')
    
    # Convert currency columns to numeric early, handling errors
    for col in ['Product Amount (TCV) ($000s)', 'Service Amount (TCV) ($000s)', 'Expected ATR ($000s)', 'Prior ATR ($000s)']:
        if col in df.columns:
            df[col] = pd.to_numeric(
                df[col].astype(str).str.replace(r'[\$,]', '', regex=True).replace('nan', '0'),
                errors='coerce'
            ).fillna(0)
    
    filtered = df[(df['Expiration Date'] >= fy_start) & (df['Expiration Date'] <= fy_end)].copy()
    
    if filtered.empty:
        print("Warning: No renewal opportunities found within the selected date range.")
        return
    
    # CREATE TOTAL ATR MAPPING BEFORE ANY FILTERING
    # This contains the complete aggregated ATR for each Deal Id (Product + Service)
    total_atr_by_deal = filtered.groupby('Deal Id')['Expected ATR ($000s)'].sum().to_dict()
    
    # Apply minimum ATR filter BEFORE selecting columns
    if min_atr > 0:
        # Calculate aggregated ATR by Deal Id across ALL data (not just filtered columns)
        deal_atr_sum = filtered.groupby('Deal Id')['Expected ATR ($000s)'].sum()
        valid_deals = deal_atr_sum[deal_atr_sum >= min_atr].index
        
        # Filter to keep only valid deals
        filtered = filtered[filtered['Deal Id'].isin(valid_deals)]
        
        if filtered.empty:
            print(f"Warning: No opportunities found with aggregated Expected ATR >= ${int(min_atr)}K")
            return
        
        print(f"Info: Filtered to {len(valid_deals)} deals with aggregated ATR >= ${int(min_atr)}K")
    
    # NOW select only the columns we need
    filtered = filtered[BASE_SLIDE_COLUMNS + ['Customer Name']]
    
    file_creation_date = get_file_creation_date(excel_filename)

    # Build filename suffix with MIN_ATR if applicable
    fy_suffix = f"_{initial_fy}-{final_fy}"
    if min_atr > 0:
        atr_suffix = f"_MIN_ATR_{int(min_atr)}K"
        fy_suffix = fy_suffix + atr_suffix
    
    # For product and service filenames
    product_pptx_filename_suffix = f"_product{fy_suffix}"
    service_pptx_filename_suffix = f"_service{fy_suffix}"

    # --- PRODUCT PPT ---
    filtered_product = filtered[filtered['Product Amount (TCV) ($000s)'] > 0]
    product_columns = [c for c in BASE_SLIDE_COLUMNS if c != 'Service Amount (TCV) ($000s)']
    if not filtered_product.empty:
        create_ppt_for_filter(
            filtered_product, fy_start, fy_end, excel_filename, product_pptx_filename_suffix,
            initial_fy, final_fy, file_creation_date,
            "Product Renewal", product_columns, min_atr, total_atr_by_deal, template_pptx
        )
    else:
        print("No product opportunities found.")

    # --- SERVICE PPT ---
    filtered_service = filtered[filtered['Service Amount (TCV) ($000s)'] > 0]
    service_columns = [c for c in BASE_SLIDE_COLUMNS if c != 'Product Amount (TCV) ($000s)']
    if not filtered_service.empty:
        create_ppt_for_filter(
            filtered_service, fy_start, fy_end, excel_filename, service_pptx_filename_suffix,
            initial_fy, final_fy, file_creation_date,
            "Service Renewal", service_columns, min_atr, total_atr_by_deal, template_pptx
        )
    else:
        print("No service opportunities found.")

if __name__ == '__main__':
    parser = argparse.ArgumentParser(description="Generate Cisco renewal opportunity PowerPoints for products and services from Excel.")
    parser.add_argument('initial_fy', type=str, help="Initial date (format: QQFYXX, e.g., Q1FY26)")
    parser.add_argument('final_fy', type=str, help="Final date (format: QQFYXX, e.g., Q3FY26)")
    parser.add_argument('excel_filename', type=str, help="Input Excel file name (.xlsx)")
    parser.add_argument('--min-atr', type=float, dest='min_atr', default=0, help="Minimum Expected ATR ($000s) threshold (optional, default: 0)")
    parser.add_argument('--template-pptx', type=str, dest='template_pptx', default=None, help="Optional PowerPoint template file (.pptx)")
    args = parser.parse_args()

    # Robust parameter validation
    if not robust_validate_fy_quarter(args.initial_fy):
        sys.exit(2)
    if not robust_validate_fy_quarter(args.final_fy):
        sys.exit(2)
    if not robust_check_excel_file(args.excel_filename):
        sys.exit(2)
    if not robust_check_template_file(args.template_pptx):
        sys.exit(2)

    try:
        create_renewal_ppt(args.initial_fy, args.final_fy, args.excel_filename, args.min_atr, args.template_pptx)
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        import traceback
        traceback.print_exc()
        sys.exit(2)

