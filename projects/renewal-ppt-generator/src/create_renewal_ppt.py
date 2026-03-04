'''
Baseline renewals PPT generator.

Generates Cisco renewal opportunity PowerPoints for product and service renewals
from CS Console Excel exports.

Usage:
    python create_renewal_ppt.py <initial_fy> <final_fy> <excel_filename>

Example:
    python create_renewal_ppt.py Q1FY26 Q3FY26 .\\data\\renewals\\renewals.xlsx

Inputs:
    initial_fy      Fiscal quarter start (QQFYXX), e.g. Q1FY26
    final_fy        Fiscal quarter end (QQFYXX), e.g. Q3FY26
    excel_filename  Renewals export file (.xlsx)

Outputs:
    - <input>_product_<FY-range>.pptx
    - <input>_service_<FY-range>.pptx

Note:
    For richer multi-customer summaries and threshold filtering, use
    `create_renew_ops_ppt.py`.
'''


import argparse
import io
import os
import sys
from datetime import datetime

import matplotlib.lines as mlines
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import pandas as pd
from matplotlib.dates import DateFormatter, MonthLocator
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

# --- CONSTANTS ---
CISCO_FY_QUARTERS = {
    'Q1': (8, 1),   # August 1
    'Q2': (11, 1),  # November 1
    'Q3': (2, 1),   # February 1
    'Q4': (5, 1)    # May 1
}
REQUIRED_COLUMNS = [
    'Account ARR ($000s)', 'Account Name', 'CX Upsell/PMG', 'Close Date', 'Customer Name',
    'Customer Pulse', 'Deal Id', 'Deal Pulse', 'Expected ATR ($000s)', 'Expiration Date',
    'Expiration Quarter', 'Linked/Related', 'Linked/Related Deals', 'Opportunity Name',
    'Opportunity Owner', 'Opportunity Status', 'Prior ATR ($000s)', 'Product Amount (TCV) ($000s)',
    'Renewal Risk', 'Service Amount (TCV) ($000s)', 'Stage', 'Success Priority'
]
COLUMN_ALIASES = {
    'Renewal Risk': 'RenewalLine Risk'
}
BASE_SLIDE_COLUMNS = [
    'Account Name', 'Deal Id', 'Opportunity Name', 'CX Upsell/PMG',
    'Prior ATR ($000s)', 'Expected ATR ($000s)',
    'Product Amount (TCV) ($000s)',
    'Service Amount (TCV) ($000s)', 'Expiration Date',
    'Deal Pulse', 'Customer Pulse'  # Needed for color/suffix logic
]
MAX_TABLE_ROWS_PER_SLIDE = 11
MAX_TIMELINES_PER_SLIDE = 16
TIMELINE_SPACING = 0.4  # inches between timelines
BASE_HEIGHT = 1.5       # top/bottom padding for figure in inches

MIN_CIRCLE_SIZE = 4
MAX_CIRCLE_SIZE = 16

# Color name to RGB mapping for python-pptx and matplotlib
COLOR_RGB = {
    'black': RGBColor(0, 0, 0),
    'red': RGBColor(255, 0, 0),
    'green': RGBColor(0, 128, 0),
    'yellow': RGBColor(255, 215, 0),
    'blue': RGBColor(0, 0, 255),
    'orange': RGBColor(255, 140, 0),
    'purple': RGBColor(128, 0, 128),
    'grey': RGBColor(128, 128, 128),
}
DEFAULT_COLOR = 'black'


def robust_validate_fy_quarter(fy_str):
    if not isinstance(fy_str, str):
        print(f"Error: Fiscal quarter must be a string, got {type(fy_str)}", file=sys.stderr)
        return False
    if len(fy_str) != 6:
        print(f"Error: Fiscal quarter '{fy_str}' must have 6 characters (e.g., Q1FY26).", file=sys.stderr)
        return False
    quarter = fy_str[:2]
    if quarter not in CISCO_FY_QUARTERS:
        print(f"Error: Quarter '{quarter}' in '{fy_str}' is invalid. Expected one of {list(CISCO_FY_QUARTERS)}.", file=sys.stderr)
        return False
    if fy_str[2:4] != 'FY':
        print(f"Error: Fiscal quarter '{fy_str}' must contain 'FY' after the quarter.", file=sys.stderr)
        return False
    if not fy_str[4:].isdigit():
        print(f"Error: Year '{fy_str[4:]}' in '{fy_str}' is not numeric.", file=sys.stderr)
        return False
    return True


def robust_check_excel_file(filename):
    if not isinstance(filename, str):
        print('Error: Excel filename must be a string.', file=sys.stderr)
        return False
    if not filename.lower().endswith('.xlsx'):
        print("Error: Excel filename must end with '.xlsx'.", file=sys.stderr)
        return False
    if not os.path.isfile(filename):
        print(f"Error: File '{filename}' does not exist.", file=sys.stderr)
        return False
    return True


def validate_fy_quarter(fy_str):
    quarter = fy_str[:2]
    year = int('20' + fy_str[-2:])
    month, day = CISCO_FY_QUARTERS[quarter]
    if quarter == 'Q1':
        start = datetime(year - 1, month, day)
        end = datetime(year - 1, 10, 31)
    elif quarter == 'Q2':
        start = datetime(year - 1, month, day)
        end = datetime(year, 1, 31)
    elif quarter == 'Q3':
        start = datetime(year, month, day)
        end = datetime(year, 4, 30)
    elif quarter == 'Q4':
        start = datetime(year, month, day)
        end = datetime(year, 7, 31)
    else:
        raise ValueError(f'Unknown quarter: {quarter}')
    return start, end


def get_fy_range(start_fy, end_fy):
    start_dt, _ = validate_fy_quarter(start_fy)
    _, end_dt = validate_fy_quarter(end_fy)
    return start_dt, end_dt


def get_file_creation_date(path):
    t = os.path.getctime(path)
    return datetime.fromtimestamp(t).strftime('%Y-%m-%d')


def check_excel_file(filename):
    if not os.path.isfile(filename):
        raise FileNotFoundError(f"Input Excel file '{filename}' does not exist.")
    try:
        df = pd.read_excel(filename)
    except Exception as e:
        raise ValueError(f'Could not read Excel file: {e}')
    for expected, actual in COLUMN_ALIASES.items():
        if expected not in df.columns and actual in df.columns:
            print(f"Info: Using '{actual}' as '{expected}'")
            df[expected] = df[actual]
    missing = [col for col in REQUIRED_COLUMNS if col not in df.columns]
    if missing:
        raise ValueError(f'Excel file missing required columns: {missing}')
    return df


def get_pulse_color(pulse_val):
    """
    Extract color name from Deal Pulse or Customer Pulse.
    Accepts values like '1 - Yellow', '2- yellow ', 'NA'.
    Returns color name as lower case, or 'black' if not found.
    """
    if isinstance(pulse_val, str) and pulse_val.strip().upper() != 'NA':
        parts = pulse_val.split('-')
        if len(parts) == 2:
            color_name = parts[1].strip().lower()
            return color_name if color_name in COLOR_RGB else DEFAULT_COLOR
    return DEFAULT_COLOR


def add_summary_table_slide(prs, filtered, fy_start, fy_end, value_col, title_prefix):
    months = pd.date_range(start=fy_start, end=fy_end, freq='MS')
    month_labels = [m.strftime('%b %Y') for m in months]
    account_names = sorted(filtered['Account Name'].unique())
    filtered = filtered.copy()
    filtered['Month'] = filtered['Expiration Date'].dt.to_period('M')
    table_data = pd.DataFrame(0.0, index=account_names, columns=month_labels)
    for account in account_names:
        acc_df = filtered[filtered['Account Name'] == account]
        for m, label in zip(months, month_labels):
            month_val = acc_df[(acc_df['Expiration Date'].dt.month == m.month) & (acc_df['Expiration Date'].dt.year == m.year)]
            total = month_val[value_col].sum()
            table_data.loc[account, label] = total

    table_data['Total ($000s)'] = table_data.sum(axis=1)
    total_row = table_data.sum(axis=0)
    table_cols = list(month_labels) + ['Total ($000s)']

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    for shape in list(slide.shapes):
        if shape.is_placeholder and shape.placeholder_format.type == 1:
            sp = shape
            sp.element.getparent().remove(sp.element)
    left = Inches(0.2)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.5)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.text = f'{title_prefix} Summary Table - Expected ATR Aggregated by Month ($000s)'
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].alignment = 1

    rows = table_data.shape[0]
    cols = len(table_cols)
    table_height = Inches(0.4 + 0.22 * min(rows, 40))
    table = slide.shapes.add_table(rows + 2, cols + 1, Inches(0.2), Inches(1.0), Inches(9), table_height).table
    table.cell(0, 0).text = 'Account Name'
    table.cell(0, 0).text_frame.paragraphs[0].font.size = Pt(8)
    table.cell(0, 0).text_frame.paragraphs[0].font.bold = True
    for j, label in enumerate(table_cols):
        cell = table.cell(0, j + 1)
        cell.text = label
        cell.text_frame.paragraphs[0].font.size = Pt(8)
        cell.text_frame.paragraphs[0].font.bold = True

    for i, account in enumerate(table_data.index):
        table.cell(i + 1, 0).text = account
        table.cell(i + 1, 0).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(i + 1, 0).text_frame.paragraphs[0].font.bold = True
        for j, label in enumerate(table_cols):
            val = table_data.loc[account, label]
            cell = table.cell(i + 1, j + 1)
            cell.text = f'${int(round(val)):,}' if val > 0 else ''
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            if label == 'Total ($000s)':
                cell.text_frame.paragraphs[0].font.bold = True

    table.cell(rows + 1, 0).text = 'Total ($000s)'
    table.cell(rows + 1, 0).text_frame.paragraphs[0].font.size = Pt(8)
    table.cell(rows + 1, 0).text_frame.paragraphs[0].font.bold = True
    for j, label in enumerate(table_cols):
        val = total_row[label]
        cell = table.cell(rows + 1, j + 1)
        cell.text = f'${int(round(val)):,}' if val > 0 else ''
        cell.text_frame.paragraphs[0].font.size = Pt(8)
        cell.text_frame.paragraphs[0].font.bold = True
    return slide


def add_table_slides(prs, filtered, columns, title_prefix, account_col='Account Name'):
    for account in filtered[account_col].unique():
        acc_table_df = filtered[filtered[account_col] == account][columns]
        num_rows = acc_table_df.shape[0]
        num_slides = (num_rows // MAX_TABLE_ROWS_PER_SLIDE) + (1 if num_rows % MAX_TABLE_ROWS_PER_SLIDE else 0)
        for slide_idx, start_row in enumerate(range(0, num_rows, MAX_TABLE_ROWS_PER_SLIDE)):
            chunk = acc_table_df.iloc[start_row:start_row + MAX_TABLE_ROWS_PER_SLIDE]
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            for shape in list(slide.shapes):
                if shape.is_placeholder and shape.placeholder_format.type == 1:
                    sp = shape
                    sp.element.getparent().remove(sp.element)
            left = Inches(0.2)
            top = Inches(0.3)
            width = Inches(9)
            height = Inches(0.5)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            tf.text = f'{account} {title_prefix} Opportunities' + (f' (page {slide_idx + 1})' if num_slides > 1 else '')
            tf.paragraphs[0].font.size = Pt(16)
            tf.paragraphs[0].alignment = 1
            rows, cols = chunk.shape
            table_height = Inches(0.4 + 0.25 * min(rows, MAX_TABLE_ROWS_PER_SLIDE))
            table = slide.shapes.add_table(rows + 1, cols, Inches(0.2), Inches(1.0), Inches(9), table_height).table
            for j, col in enumerate(chunk.columns):
                cell = table.cell(0, j)
                cell.text = str(col)
                cell.text_frame.paragraphs[0].font.size = Pt(8)
                cell.text_frame.paragraphs[0].font.bold = True
            for i, row in enumerate(chunk.values):
                for j, _colname in enumerate(chunk.columns):
                    cell = table.cell(i + 1, j)
                    val = '' if pd.isnull(row[j]) else row[j].strftime('%Y-%m-%d') if isinstance(row[j], pd.Timestamp) else str(row[j])
                    cell.text = val
                    cell.text_frame.paragraphs[0].font.size = Pt(6)


def add_timeline_legend(ax):
    pulse_colors = {
        'green': '#34A853',
        'yellow': '#FBBC05',
        'red': '#EA4335',
        'grey': '#B0B0B0'
    }
    legend_elements = [
        mlines.Line2D([], [], color=pulse_colors['green'], marker='o', linestyle='None', markersize=9, label='Deal Pulse color & ATR size in circle'),
        mpatches.Patch(facecolor='white', edgecolor='gray', label='Customer Pulse color in Deal Id text', alpha=0.0)
    ]
    ax.legend(
        handles=legend_elements,
        loc='lower center',
        bbox_to_anchor=(0.5, 1.02),
        ncol=2,
        fontsize=9,
        frameon=True
    )


def get_circle_size(atr, min_atr, max_atr, min_circle_size=MIN_CIRCLE_SIZE, max_circle_size=MAX_CIRCLE_SIZE):
    if max_atr == min_atr:
        return min_circle_size
    norm = (atr - min_atr) / (max_atr - min_atr)
    return min_circle_size + norm * (max_circle_size - min_circle_size)


def get_atr_min_max(df, account_name, atr_column='Expected ATR ($000s)'):
    deals = df[df['Account Name'] == account_name]
    atr_values = pd.to_numeric(deals[atr_column], errors='coerce').dropna()
    if atr_values.empty:
        return 1.0, 1.0
    min_atr = max(atr_values.min(), 1)
    max_atr = atr_values.max()
    return min_atr, max_atr


def add_timeline_slides(prs, filtered, fy_start, fy_end, title_prefix, account_col='Account Name'):
    for account in filtered[account_col].unique():
        min_atr, max_atr = get_atr_min_max(filtered, account)
        acc_df = filtered[filtered[account_col] == account].copy()
        acc_df = acc_df.sort_values('Expiration Date').reset_index(drop=True)
        timelines = []
        timeline_assignments = []
        for _idx, row in acc_df.iterrows():
            placed = False
            for t_idx, timeline in enumerate(timelines):
                if (row['Expiration Date'] - timeline[-1]['Expiration Date']).days >= 15:
                    timeline.append(row)
                    timeline_assignments.append(t_idx)
                    placed = True
                    break
            if not placed:
                timelines.append([row])
                timeline_assignments.append(len(timelines) - 1)
        acc_df['timeline_row'] = timeline_assignments
        n_timelines = len(timelines)
        for slide_idx, start_timeline in enumerate(range(0, n_timelines, MAX_TIMELINES_PER_SLIDE)):
            end_timeline = start_timeline + MAX_TIMELINES_PER_SLIDE
            acc_slide_df = acc_df[(acc_df['timeline_row'] >= start_timeline) & (acc_df['timeline_row'] < end_timeline)].copy()
            if acc_slide_df.empty:
                continue
            min_row = acc_slide_df['timeline_row'].min()
            acc_slide_df['plot_row'] = acc_slide_df['timeline_row'] - min_row
            n_rows_this_slide = acc_slide_df['plot_row'].max() + 1
            fig_height = BASE_HEIGHT + TIMELINE_SPACING * n_rows_this_slide
            fig, ax = plt.subplots(figsize=(10, fig_height))

            for i in range(n_rows_this_slide):
                ax.hlines(y=i, xmin=fy_start, xmax=fy_end, color='tab:blue', linewidth=2)

            for _, row in acc_slide_df.iterrows():
                x = row['Expiration Date']
                y = row['plot_row']
                deal_id = str(row['Deal Id'])
                deal_pulse_color = get_pulse_color(row.get('Deal Pulse', 'NA')).lower()
                customer_pulse_color = get_pulse_color(row.get('Customer Pulse', 'NA')).lower()
                atr = float(row['Expected ATR ($000s)'])
                circle_size = get_circle_size(atr, min_atr, max_atr)
                ax.plot(x, y, 'o', color=deal_pulse_color, markersize=circle_size)

                if customer_pulse_color == 'yellow':
                    ax.text(
                        x, y + 0.12, deal_id,
                        fontsize=8, ha='center', va='bottom', rotation=20, fontweight='bold',
                        color='black',
                        bbox=dict(facecolor='yellow', edgecolor='none', boxstyle='round,pad=0.1')
                    )
                else:
                    mpl_text_color = customer_pulse_color if customer_pulse_color in COLOR_RGB else 'black'
                    if mpl_text_color == 'yellow':
                        mpl_text_color = 'black'
                    ax.text(
                        x, y + 0.12, deal_id,
                        fontsize=8, ha='center', va='bottom', rotation=20, fontweight='bold',
                        color=mpl_text_color
                    )

            months = pd.date_range(start=fy_start, end=fy_end, freq='MS')
            for month_start in months:
                ax.vlines(month_start, ymin=-0.5, ymax=n_rows_this_slide - 0.5, color='lightgrey', linestyle='--', linewidth=1)
            ax.set_yticks([])
            ax.set_xlim([fy_start, fy_end])
            ax.set_ylim(-0.5, n_rows_this_slide - 0.5)
            ax.set_title('')
            ax.set_xlabel('Expiration Date')
            ax.xaxis.set_major_locator(MonthLocator())
            ax.xaxis.set_major_formatter(DateFormatter('%b\n%Y'))

            add_timeline_legend(ax)

            plt.tight_layout()
            img_stream = io.BytesIO()
            plt.savefig(img_stream, format='png')
            plt.close(fig)
            img_stream.seek(0)
            timeline_slide = prs.slides.add_slide(prs.slide_layouts[5])
            for shape in list(timeline_slide.shapes):
                if shape.is_placeholder and shape.placeholder_format.type == 1:
                    sp = shape
                    sp.element.getparent().remove(sp.element)
            left = Inches(1)
            top = Inches(0.3)
            width = Inches(8)
            height = Inches(0.5)
            textbox = timeline_slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            tf.text = f'{account} {title_prefix} Opportunities Timeline' + (f' (page {slide_idx + 1})' if n_timelines > MAX_TIMELINES_PER_SLIDE else '')
            tf.paragraphs[0].font.size = Pt(16)
            tf.paragraphs[0].alignment = 1
            timeline_slide.shapes.add_picture(img_stream, Inches(1), Inches(1.0), width=Inches(8))

            month_deals = {}
            for _, row in acc_slide_df.iterrows():
                exp_date = row['Expiration Date']
                month_str = exp_date.strftime('%b %Y')
                deal_id = str(row['Deal Id'])
                atr = row.get('Expected ATR ($000s)', '')
                opp_name = row.get('Opportunity Name', '')
                upsell = row.get('CX Upsell/PMG', '')
                detail = f'{deal_id}: {atr} | {opp_name} | {upsell}'
                if month_str not in month_deals:
                    month_deals[month_str] = []
                month_deals[month_str].append(detail)
            notes_lines = []
            for month in sorted(month_deals, key=lambda x: datetime.strptime(x, '%b %Y')):
                notes_lines.append(f'{month}:')
                notes_lines.extend([f'- {deal_detail}' for deal_detail in month_deals[month]])

            notes_text = '\n'.join(notes_lines)
            timeline_slide.notes_slide.notes_text_frame.text = notes_text


def create_ppt_for_filter(filtered, fy_start, fy_end, excel_filename, suffix, customer_name, initial_fy, final_fy, file_creation_date, title_prefix, columns):
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    title.text = f'{title_prefix} Opportunities in {customer_name}'
    subtitle.text = (
        f'Date Range: {initial_fy} - {final_fy}\n'
        f'Fiscal Dates: {fy_start.strftime("%b %d, %Y")} - {fy_end.strftime("%b %d, %Y")}\n'
        f'File Created: {file_creation_date}'
    )

    value_col = 'Expected ATR ($000s)'
    add_summary_table_slide(prs, filtered, fy_start, fy_end, value_col, title_prefix)

    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.insert(1, slides[-1])
    del slides[-1]

    add_table_slides(prs, filtered, columns, title_prefix)
    add_timeline_slides(prs, filtered, fy_start, fy_end, title_prefix)
    pptx_filename = os.path.splitext(excel_filename)[0] + suffix + '.pptx'
    prs.save(pptx_filename)
    print(f'Presentation saved as: {pptx_filename}')


def create_renewal_ppt(initial_fy, final_fy, excel_filename):
    fy_start, fy_end = get_fy_range(initial_fy, final_fy)
    df = check_excel_file(excel_filename)
    if len(df) == 0:
        print('Warning: The input Excel file contains no data.')
    df['Expiration Date'] = pd.to_datetime(df['Expiration Date'], errors='coerce')
    filtered = df[(df['Expiration Date'] >= fy_start) & (df['Expiration Date'] <= fy_end)].copy()
    if filtered.empty:
        print('Warning: No renewal opportunities found within the selected date range.')
    filtered = filtered[BASE_SLIDE_COLUMNS]
    for col in ['Product Amount (TCV) ($000s)', 'Service Amount (TCV) ($000s)', 'Expected ATR ($000s)']:
        filtered[col] = (
            filtered[col]
            .astype(str)
            .str.replace(r'[\$,]', '', regex=True)
            .replace('nan', '0')
            .astype(float)
        )
    customer_name = df['Customer Name'].iloc[0] if not df.empty and 'Customer Name' in df else ''
    file_creation_date = get_file_creation_date(excel_filename)

    fy_suffix = f'_{initial_fy}-{final_fy}'
    product_pptx_filename_suffix = f'_product{fy_suffix}'
    service_pptx_filename_suffix = f'_service{fy_suffix}'

    filtered_product = filtered[filtered['Product Amount (TCV) ($000s)'] > 0]
    product_columns = [c for c in BASE_SLIDE_COLUMNS if c != 'Service Amount (TCV) ($000s)']
    if not filtered_product.empty:
        create_ppt_for_filter(
            filtered_product, fy_start, fy_end, excel_filename, product_pptx_filename_suffix,
            customer_name, initial_fy, final_fy, file_creation_date,
            'Product Renewal', product_columns
        )
    else:
        print('No product opportunities found.')

    filtered_service = filtered[filtered['Service Amount (TCV) ($000s)'] > 0]
    service_columns = [c for c in BASE_SLIDE_COLUMNS if c != 'Product Amount (TCV) ($000s)']
    if not filtered_service.empty:
        create_ppt_for_filter(
            filtered_service, fy_start, fy_end, excel_filename, service_pptx_filename_suffix,
            customer_name, initial_fy, final_fy, file_creation_date,
            'Service Renewal', service_columns
        )
    else:
        print('No service opportunities found.')


if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='Generate Cisco renewal opportunity PowerPoints for products and services from Excel.')
    parser.add_argument('initial_fy', type=str, help='Initial date (format: QQFYXX, e.g., Q1FY26)')
    parser.add_argument('final_fy', type=str, help='Final date (format: QQFYXX, e.g., Q3FY26)')
    parser.add_argument('excel_filename', type=str, help='Input Excel file name (.xlsx)')
    args = parser.parse_args()

    if not robust_validate_fy_quarter(args.initial_fy):
        sys.exit(2)
    if not robust_validate_fy_quarter(args.final_fy):
        sys.exit(2)
    if not robust_check_excel_file(args.excel_filename):
        sys.exit(2)

    try:
        create_renewal_ppt(args.initial_fy, args.final_fy, args.excel_filename)
    except Exception as e:
        print(f'Error: {e}', file=sys.stderr)
        sys.exit(2)
