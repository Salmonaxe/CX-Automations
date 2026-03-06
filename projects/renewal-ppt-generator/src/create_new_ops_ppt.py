'''
New opportunities PPT generator.

Generates Cisco new-opportunity PowerPoints from CS Console Excel exports.

Usage:
    python create_new_ops_ppt.py <initial_fy> <final_fy> <excel_filename> [--min-tcv <value>] [--template-pptx <file>]

Example:
    python create_new_ops_ppt.py Q1FY26 Q3FY26 .\\data\\new-ops\\new_ops.xlsx --min-tcv 100

Inputs:
    initial_fy      Fiscal quarter start (QQFYXX), e.g. Q1FY26
    final_fy        Fiscal quarter end (QQFYXX), e.g. Q3FY26
    excel_filename  New opportunities export file (.xlsx)
    --min-tcv       Optional minimum aggregated Expected TCV ($000s)
    --template-pptx Optional PowerPoint template (.pptx)

Output:
    - <input>_<FY-range>_TCV_MIN_<value>.pptx

Behavior:
    - Supports single-customer or all-customer input files
    - Aggregates rows by Deal Id
    - Adds summary slides for All Customers and per-customer views
    - Uses stage-based colors in timeline visualizations
'''


import argparse
import os
import sys
from datetime import datetime
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
from matplotlib.dates import MonthLocator, DateFormatter
import matplotlib.patches as mpatches
import matplotlib.lines as mlines
import numpy as np
import io

# --- CONSTANTS ---
CISCO_FY_QUARTERS = {
    'Q1': (8, 1),   # August 1
    'Q2': (11, 1),  # November 1
    'Q3': (2, 1),   # February 1
    'Q4': (5, 1)    # May 1
}
REQUIRED_COLUMNS = [
    'Account Name', 'CX Upsell/PMG', 'Close Date', 'Customer Name',
    'Deal Id', 'Expected Amount TCV ($000s)',
    'Linked/Related', 'Linked/Related Deals', 'Opportunity Name',
    'Opportunity Owner', 'Opportunity Status', 'Stage'
]
COLUMN_ALIASES = {
    'Expected Amount TCV ($000s)': 'Product Amount (TCV) ($000s)',
    'Deal Id': 'Deal ID',
}
BASE_SLIDE_COLUMNS = [
    'Account Name', 'Deal Id', 'Opportunity Name', 'CX Upsell/PMG',
    'Expected Amount TCV ($000s)', 'Opportunity Owner',
    'Opportunity Status', 'Stage',  # Needed for color logic
    'Close Date'
]
MAX_TABLE_ROWS_PER_SLIDE = 16
MAX_TIMELINES_PER_SLIDE = 16
TIMELINE_SPACING = 0.4  # inches between timelines
BASE_HEIGHT = 1.5       # top/bottom padding for figure in inches

# Set your base circle size (e.g., current size used for $1 ATR)
MIN_CIRCLE_SIZE = 4
MAX_CIRCLE_SIZE = 16  # or MIN_CIRCLE_SIZE * 4

# Color name to RGB mapping for python-pptx and matplotlib
COLOR_RGB = {
    'black': RGBColor(0,0,0),
    'red': RGBColor(255,0,0),
    'green': RGBColor(0,128,0),
    'yellow': RGBColor(255,215,0),
    'blue': RGBColor(0,0,255),
    'orange': RGBColor(255,140,0),
    'cyan': RGBColor(0, 255, 255),
    'grey': RGBColor(128,128,128)
}
DEFAULT_COLOR = 'black'

EMU_PER_INCH = 914400
SLIDE_MARGIN_IN = 0.35
CONTENT_BOTTOM_MARGIN_IN = 0.3
CONTENT_TITLE_HEIGHT_IN = 0.55
CONTENT_TITLE_GAP_IN = 0.12
TABLE_HEADER_HEIGHT_IN = 0.30
TABLE_ROW_HEIGHT_IN = 0.24


def remove_all_placeholders(slide):
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            sp = shape
            sp.element.getparent().remove(sp.element)


def get_slide_content_bounds(prs):
    left = Inches(SLIDE_MARGIN_IN)
    top = Inches(SLIDE_MARGIN_IN)
    right = Inches(SLIDE_MARGIN_IN)
    bottom = Inches(CONTENT_BOTTOM_MARGIN_IN)
    width = max(int(prs.slide_width - left - right), Inches(4))
    height = max(int(prs.slide_height - top - bottom), Inches(2))
    return left, top, width, height


def get_content_body_bounds(prs):
    left, top, width, content_height = get_slide_content_bounds(prs)
    title_height = Inches(CONTENT_TITLE_HEIGHT_IN)
    title_gap = Inches(CONTENT_TITLE_GAP_IN)
    body_top = top + title_height + title_gap
    body_height = max(int(content_height - title_height - title_gap), Inches(1.6))
    return left, top, width, title_height, body_top, body_height


def add_content_title(slide, prs, text, font_size=16):
    left, top, width, title_height, body_top, body_height = get_content_body_bounds(prs)
    textbox = slide.shapes.add_textbox(left, top, width, title_height)
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.text = text
    tf.paragraphs[0].font.size = Pt(font_size)
    tf.paragraphs[0].alignment = 1
    return left, body_top, width, body_height


def get_table_rows_per_slide(content_height_emu, include_totals=False):
    content_height_in = content_height_emu / EMU_PER_INCH
    reserved_in = TABLE_HEADER_HEIGHT_IN + (TABLE_ROW_HEIGHT_IN if include_totals else 0.0)
    usable_in = max(content_height_in - reserved_in, TABLE_ROW_HEIGHT_IN)
    return max(1, int(usable_in / TABLE_ROW_HEIGHT_IN))


def get_table_font_sizes(col_count):
    if col_count >= 10:
        return Pt(7), Pt(6)
    if col_count >= 8:
        return Pt(8), Pt(6)
    return Pt(9), Pt(7)


def apply_table_dimensions(table, total_width_emu, total_height_emu, first_col_ratio=0.30):
    col_count = len(table.columns)
    row_count = len(table.rows)
    if col_count > 1:
        first_col_width = int(total_width_emu * first_col_ratio)
        remaining = max(total_width_emu - first_col_width, 0)
        other_col_width = int(remaining / (col_count - 1)) if col_count > 1 else remaining
        table.columns[0].width = first_col_width
        for idx in range(1, col_count):
            table.columns[idx].width = other_col_width
    row_height = max(int(total_height_emu / max(row_count, 1)), Inches(0.18))
    for row in table.rows:
        row.height = row_height


def add_fitted_picture(slide, img_stream, left, top, max_width, max_height, image_aspect_ratio):
    if max_width <= 0 or max_height <= 0:
        return
    image_aspect = image_aspect_ratio if image_aspect_ratio and image_aspect_ratio > 0 else (16.0 / 9.0)
    box_aspect = max_width / max_height
    if image_aspect >= box_aspect:
        draw_width = max_width
        draw_height = int(draw_width / image_aspect)
    else:
        draw_height = max_height
        draw_width = int(draw_height * image_aspect)
    draw_left = left + int((max_width - draw_width) / 2)
    draw_top = top + int((max_height - draw_height) / 2)
    slide.shapes.add_picture(img_stream, draw_left, draw_top, width=draw_width, height=draw_height)

def robust_validate_fy_quarter(fy_str):
    if not isinstance(fy_str, str):
        print(f"Error: Fiscal quarter must be a string, got {type(fy_str)}", file=sys.stderr)
        return False
    if len(fy_str) != 6:
        print(f"Error: Fiscal quarter '{fy_str}' must have 6 characters (e.g., Q1FY26).", file=sys.stderr)
        return False
    quarter = fy_str[:2]
    if quarter not in CISCO_FY_QUARTERS:
        print(f"Error: Quarter '{quarter}' in '{fy_str}' is invalid. Expected one of {list(CISCO_FY_QUARTERS)}.", file=sys.stderr)
        return False
    if fy_str[2:4] != 'FY':
        print(f"Error: Fiscal quarter '{fy_str}' must contain 'FY' after the quarter.", file=sys.stderr)
        return False
    if not fy_str[4:].isdigit():
        print(f"Error: Year '{fy_str[4:]}' in '{fy_str}' is not numeric.", file=sys.stderr)
        return False
    return True

def robust_check_excel_file(filename):
    if not isinstance(filename, str):
        print("Error: Excel filename must be a string.", file=sys.stderr)
        return False
    if not filename.lower().endswith('.xlsx'):
        print("Error: Excel filename must end with '.xlsx'.", file=sys.stderr)
        return False
    if not os.path.isfile(filename):
        print(f"Error: File '{filename}' does not exist.", file=sys.stderr)
        return False
    return True

def robust_check_template_file(filename):
    if filename is None:
        return True
    if not isinstance(filename, str):
        print("Error: Template filename must be a string.", file=sys.stderr)
        return False
    if not filename.lower().endswith('.pptx'):
        print("Error: Template file must end with '.pptx'. If you have a .potx, save/convert it to .pptx first.", file=sys.stderr)
        return False
    if not os.path.isfile(filename):
        print(f"Error: Template file '{filename}' does not exist.", file=sys.stderr)
        return False
    return True

def validate_fy_quarter(fy_str):
    quarter = fy_str[:2]
    year = int('20' + fy_str[-2:])
    month, day = CISCO_FY_QUARTERS[quarter]
    if quarter == 'Q1':
        start = datetime(year - 1, month, day)
        end = datetime(year - 1, 10, 31)
    elif quarter == 'Q2':
        start = datetime(year - 1, month, day)
        end = datetime(year, 1, 31)
    elif quarter == 'Q3':
        start = datetime(year, month, day)
        end = datetime(year, 4, 30)
    elif quarter == 'Q4':
        start = datetime(year, month, day)
        end = datetime(year, 7, 31)
    else:
        raise ValueError(f"Unknown quarter: {quarter}")
    return start, end

def get_fy_range(start_fy, end_fy):
    start_dt, _ = validate_fy_quarter(start_fy)
    _, end_dt = validate_fy_quarter(end_fy)
    return start_dt, end_dt

def get_file_creation_date(path):
    t = os.path.getctime(path)
    return datetime.fromtimestamp(t).strftime('%Y-%m-%d')

def check_excel_file(filename):
    if not os.path.isfile(filename):
        raise FileNotFoundError(f"Input Excel file '{filename}' does not exist.")
    try:
        df = pd.read_excel(filename)
    except Exception as e:
        raise ValueError(f"Could not read Excel file: {e}")
    for expected, actual in COLUMN_ALIASES.items():
        if expected not in df.columns and actual in df.columns:
            print(f"Info: Using '{actual}' as '{expected}'")
            df[expected] = df[actual]
    missing = [col for col in REQUIRED_COLUMNS if col not in df.columns]
    if missing:
        raise ValueError(f"Excel file missing required columns: {missing}")
    return df

def get_stage_color(stage_val):
    """
    Extract color name from Stage.
    Accepts values: 
        '1 - Qualify',              #black
        '2 - Propose',              #blue
        '3 - Technical Validation', #cyan
        '4 - Business Validation',  #yellow
        '5 - Negotiate',            #orange
        '6 - Closed Won',           #green
        '6 - Closed Lost'           #red

    Checks value in stage_val. Strips the string and gets values (Qualify, Propose, ...)
    Then returns associated color name as shown above, in lower case.
    """
    if not stage_val or not isinstance(stage_val, str):
        return None

    stage_to_color = {
        "qualify": "black",
        "propose": "blue",
        "technical validation": "cyan",
        "business validation": "yellow",
        "negotiate": "orange",
        "closed won": "green",
        "closed lost": "red",
    }

    # Strip whitespace and split on "-" to extract stage name
    try:
        stage_name = stage_val.split("-", 1)[1].strip().lower()
    except IndexError:
        return None

    return stage_to_color.get(stage_name)


def get_customer_names(df):
    if 'Customer Name' not in df.columns:
        return []
    customer_names = (
        df['Customer Name']
        .dropna()
        .astype(str)
        .str.strip()
    )
    customer_names = customer_names[customer_names != '']
    return sorted(customer_names.unique().tolist())


def format_customer_scope(customer_names, max_names=4):
    if not customer_names:
        return 'Selected Customers'
    if len(customer_names) <= max_names:
        return ', '.join(customer_names)
    return f"{', '.join(customer_names[:max_names])} (+{len(customer_names) - max_names} more)"


def add_summary_table_slide(
    prs, filtered, fy_start, fy_end, value_col, summary_scope='All Customers', opportunity_status=None
):
    """
    Adds a summary table slide to the PowerPoint presentation.

    Parameters:
        prs: pptx.Presentation object.
        filtered: pd.DataFrame, must have 'Account Name', 'Close Date', and value_col.
        fy_start: str or datetime, start of fiscal year.
        fy_end: str or datetime, end of fiscal year.
        value_col: str, column to aggregate.
        summary_scope: str, display label for the scope (e.g., all customers or one customer).
        opportunity_status: optional str in {'Lost', 'Booked', 'Active'} to limit rows.

    Returns:
        slide: pptx slide object.
    """
    # Define mapping between status and title
    status_map = {
        'Lost':    ("Lost",    "Lost Opportunities"),
        'Booked':  ("Booked",  "Booked Opportunities"),
        'Active':  ("Active",  "Active Opportunities"),
    }

    if opportunity_status is None:
        title_prefix = "All Statuses"
        data = filtered
    else:
        if opportunity_status not in status_map:
            raise ValueError("opportunity_status must be one of: 'Lost', 'Booked', 'Active'")
        status_value, title_prefix = status_map[opportunity_status]
        data = filtered[filtered['Opportunity Status'] == status_value]

    if data.empty:
        return None

    # Setup months and labels
    months = pd.date_range(start=fy_start, end=fy_end, freq='MS')
    month_labels = [m.strftime('%b %Y') for m in months]
    account_names = sorted(data['Account Name'].unique())
    data = data.copy()
    data['Month'] = data['Close Date'].dt.to_period('M')
    table_data = pd.DataFrame(0.0, index=account_names, columns=month_labels)

    for account in account_names:
        acc_df = data[data['Account Name'] == account]
        for m, label in zip(months, month_labels):
            month_val = acc_df[
                (acc_df['Close Date'].dt.month == m.month) & (acc_df['Close Date'].dt.year == m.year)
            ]
            total = month_val[value_col].sum()
            table_data.loc[account, label] = total

    # Add Total column and row
    table_data['Total ($000s)'] = table_data.sum(axis=1)
    total_row = table_data.sum(axis=0)
    table_cols = list(month_labels) + ['Total ($000s)']

    _left, _top, _width, _title_height, _body_top, body_height = get_content_body_bounds(prs)
    rows_per_slide = min(MAX_TABLE_ROWS_PER_SLIDE, get_table_rows_per_slide(body_height, include_totals=True))
    header_font, body_font = get_table_font_sizes(len(table_cols) + 1)
    total_accounts = list(table_data.index)
    num_pages = (len(total_accounts) + rows_per_slide - 1) // rows_per_slide

    last_slide = None
    for page_idx, start_row in enumerate(range(0, len(total_accounts), rows_per_slide), start=1):
        page_accounts = total_accounts[start_row:start_row + rows_per_slide]
        is_last_page = start_row + rows_per_slide >= len(total_accounts)
        title_suffix = f" (page {page_idx})" if num_pages > 1 else ""

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        remove_all_placeholders(slide)
        table_left, table_top, table_width, table_height = add_content_title(
            slide,
            prs,
            f"{title_prefix} Summary - {summary_scope} - Expected TCV Aggregated by Month ($000s){title_suffix}",
            font_size=16,
        )

        total_rows = 1 + len(page_accounts) + (1 if is_last_page else 0)
        table = slide.shapes.add_table(total_rows, len(table_cols) + 1, table_left, table_top, table_width, table_height).table
        apply_table_dimensions(table, table_width, table_height, first_col_ratio=0.28)

        table.cell(0, 0).text = "Account Name"
        table.cell(0, 0).text_frame.paragraphs[0].font.size = header_font
        table.cell(0, 0).text_frame.paragraphs[0].font.bold = True
        for j, label in enumerate(table_cols):
            cell = table.cell(0, j + 1)
            cell.text = label
            cell.text_frame.paragraphs[0].font.size = header_font
            cell.text_frame.paragraphs[0].font.bold = True

        for i, account in enumerate(page_accounts, start=1):
            table.cell(i, 0).text = account
            table.cell(i, 0).text_frame.paragraphs[0].font.size = body_font
            table.cell(i, 0).text_frame.paragraphs[0].font.bold = True
            for j, label in enumerate(table_cols):
                val = table_data.loc[account, label]
                cell = table.cell(i, j + 1)
                cell.text = f"${int(round(val)):,}" if val > 0 else ""
                cell.text_frame.paragraphs[0].font.size = body_font
                if label == 'Total ($000s)':
                    cell.text_frame.paragraphs[0].font.bold = True

        if is_last_page:
            total_idx = len(page_accounts) + 1
            table.cell(total_idx, 0).text = "Total ($000s)"
            table.cell(total_idx, 0).text_frame.paragraphs[0].font.size = body_font
            table.cell(total_idx, 0).text_frame.paragraphs[0].font.bold = True
            for j, label in enumerate(table_cols):
                val = total_row[label]
                cell = table.cell(total_idx, j + 1)
                cell.text = f"${int(round(val)):,}" if val > 0 else ""
                cell.text_frame.paragraphs[0].font.size = body_font
                cell.text_frame.paragraphs[0].font.bold = True
        last_slide = slide
    return last_slide

def add_table_slides(prs, filtered, columns, title_prefix, account_col='Account Name'):
    for account in filtered[account_col].unique():
        acc_table_df = filtered[filtered[account_col] == account][columns]
        num_rows = acc_table_df.shape[0]
        _left, _top, _width, _title_height, _body_top, body_height = get_content_body_bounds(prs)
        rows_per_slide = min(MAX_TABLE_ROWS_PER_SLIDE, get_table_rows_per_slide(body_height, include_totals=False))
        num_slides = (num_rows // rows_per_slide) + (1 if num_rows % rows_per_slide else 0)
        for slide_idx, start_row in enumerate(range(0, num_rows, rows_per_slide)):
            chunk = acc_table_df.iloc[start_row:start_row+rows_per_slide]
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            remove_all_placeholders(slide)
            table_left, table_top, table_width, table_height = add_content_title(
                slide,
                prs,
                f"{account} {title_prefix} " + (f"(page {slide_idx+1})" if num_slides > 1 else ""),
                font_size=15,
            )
            rows, cols = chunk.shape
            table = slide.shapes.add_table(rows+1, cols, table_left, table_top, table_width, table_height).table
            apply_table_dimensions(table, table_width, table_height, first_col_ratio=0.24)
            header_font, body_font = get_table_font_sizes(cols)
            for j, col in enumerate(chunk.columns):
                cell = table.cell(0, j)
                cell.text = str(col)
                cell.text_frame.paragraphs[0].font.size = header_font
                cell.text_frame.paragraphs[0].font.bold = True
            for i, row in enumerate(chunk.values):
                for j, colname in enumerate(chunk.columns):
                    cell = table.cell(i+1, j)
                    val = "" if pd.isnull(row[j]) else row[j].strftime('%Y-%m-%d') if isinstance(row[j], pd.Timestamp) else str(int(row[j])) if isinstance(row[j], (float, int)) else str(row[j])
                    cell.text = val
                    cell.text_frame.paragraphs[0].font.size = body_font

# Helper function to convert RGBColor object to Matplotlib-friendly tuple
def to_mpl_rgb(rgb_obj):
    # This assumes RGBColor objects can be indexed like (r, g, b)
    return (rgb_obj[0]/255.0, rgb_obj[1]/255.0, rgb_obj[2]/255.0)

def add_timeline_legend(ax):
    '''
    # Example colors, adjust to match your actual color mapping if needed
        '1 - Qualify',              #black
        '2 - Propose',              #blue
        '3 - Technical Validation', #cyan
        '4 - Business Validation',  #yellow
        '5 - Negotiate',            #orange
        '6 - Closed Won',           #green
        '6 - Closed Lost'           #red
    '''
    # Use Line2D for circles in the legend
    legend_elements = [
        mlines.Line2D([], [], color=to_mpl_rgb(COLOR_RGB['black']), marker='o', linestyle='None', markersize=9, label='Qualify'),
        mlines.Line2D([], [], color=to_mpl_rgb(COLOR_RGB['blue']), marker='o', linestyle='None', markersize=9, label='Propose'),
        mlines.Line2D([], [], color=to_mpl_rgb(COLOR_RGB['cyan']), marker='o', linestyle='None', markersize=9, label='Technical Validation'),
        mlines.Line2D([], [], color=to_mpl_rgb(COLOR_RGB['yellow']), marker='o', linestyle='None', markersize=9, label='Business Validation'),
        mlines.Line2D([], [], color=to_mpl_rgb(COLOR_RGB['orange']), marker='o', linestyle='None', markersize=9, label='Negotiate'),
        mlines.Line2D([], [], color=to_mpl_rgb(COLOR_RGB['green']), marker='o', linestyle='None', markersize=9, label='Won'),
        mlines.Line2D([], [], color=to_mpl_rgb(COLOR_RGB['red']), marker='o', linestyle='None', markersize=9, label='Lost')
    ]
    # Place legend below the plot, centered
    ax.legend(
        handles=legend_elements,
        loc='lower center',
        bbox_to_anchor=(0.5, 1.02),  # Adjust the -0.15 value as needed for spacing
        ncol=7,  # You can use more columns if you have many items
        fontsize=9,
        frameon=True
    )

    # Optionally, add a small text box for clarity
    #ax.figure.text(
    #    0.98, 0.05,
    #    'Legend:\nDeal Pulse color = circle\nCustomer Pulse color = deal Id text',
    #    fontsize=10, ha='right', va='bottom'
    #)

def get_circle_size(atr, min_atr, max_atr, min_circle_size=MIN_CIRCLE_SIZE, max_circle_size=MAX_CIRCLE_SIZE):
        if max_atr == min_atr:
            return min_circle_size
        norm = (atr - min_atr) / (max_atr - min_atr)
        return min_circle_size + norm * (max_circle_size - min_circle_size)
    # Usage:
    # min_atr, max_atr = get_atr_min_max(df, account_name)

def get_atr_min_max(df, account_name, atr_column='Expected Amount TCV ($000s)'):
    """
    Returns min_atr (minimum, at least 1) and max_atr for all deals of a given Account Name.
    df: pandas DataFrame with all deals.
    account_name: string, the Account Name to filter by.
    atr_column: the column name for Expected Amount TCV ($000s).
    """
    # Filter for the specified Account Name
    deals = df[df['Account Name'] == account_name]
    # Drop missing or non-numeric values
    atr_values = pd.to_numeric(deals[atr_column], errors='coerce').dropna()
    if atr_values.empty:
        # Return None or sensible defaults if no valid deals
        return 1.0, 1.0
    min_atr = max(atr_values.min(), 1)
    max_atr = atr_values.max()
    return min_atr, max_atr


def add_timeline_slides(prs, filtered, fy_start, fy_end, title_prefix, account_col='Account Name', min_tcv=0):
    for account in filtered[account_col].unique():
        acc_df = filtered[filtered[account_col] == account].copy()
        acc_df['Deal Id'] = acc_df['Deal Id'].astype(str).str.strip()
        
        # Aggregate by Deal Id for total TCV calculation
        deal_id_total = acc_df.groupby('Deal Id')['Expected Amount TCV ($000s)'].sum().to_dict()

        # Filter Deal Ids based on min_tcv
        filtered_deal_ids = {deal_id: total for deal_id, total in deal_id_total.items() if total >= min_tcv}
        acc_df = acc_df[acc_df['Deal Id'].isin(filtered_deal_ids.keys())]

        # Group by Deal Id + Close Date for timeline
        layout_df = (
            acc_df.groupby(['Deal Id', 'Close Date'])
            .agg({'Expected Amount TCV ($000s)': 'sum', 'Stage': 'first'})
            .reset_index()
            .sort_values('Close Date')
        )
        
        # Assign timeline rows
        timelines = []
        timeline_assignments = []
        for _, row in layout_df.iterrows():
            placed = False
            for t_idx, timeline in enumerate(timelines):
                if (row['Close Date'] - timeline[-1]['Close Date']).days >= 15:
                    timeline.append(row)
                    timeline_assignments.append(t_idx)
                    placed = True
                    break
            if not placed:
                timelines.append([row])
                timeline_assignments.append(len(timelines) - 1)
        layout_df['timeline_row'] = timeline_assignments
        
        # Map back to acc_df for notes
        acc_df = pd.merge(acc_df, layout_df[['Deal Id', 'Close Date', 'timeline_row']], on=['Deal Id', 'Close Date'], how='left')

        n_timelines = len(timelines)
        for slide_idx, start_timeline in enumerate(range(0, n_timelines, MAX_TIMELINES_PER_SLIDE)):
            end_timeline = start_timeline + MAX_TIMELINES_PER_SLIDE
            
            # Plotting data for this slide
            slide_plot_df = layout_df[(layout_df['timeline_row'] >= start_timeline) & (layout_df['timeline_row'] < end_timeline)].copy()
            acc_slide_df = acc_df[(acc_df['timeline_row'] >= start_timeline) & (acc_df['timeline_row'] < end_timeline)].copy()

            if slide_plot_df.empty:
                continue

            # Ensure min_atr and max_atr are defined
            unique_deal_ids_on_slide = slide_plot_df['Deal Id'].unique()
            slide_totals = [deal_id_total[deal_id] for deal_id in unique_deal_ids_on_slide]

            # Initialize min_atr and max_atr to handle empty lists
            if slide_totals:
                min_atr = min(slide_totals)
                max_atr = max(slide_totals)
            else:
                min_atr, max_atr = 1.0, 1.0  # Default values

            # Prepare the plot
            min_row = slide_plot_df['timeline_row'].min()
            slide_plot_df['plot_row'] = slide_plot_df['timeline_row'] - min_row
            n_rows_this_slide = slide_plot_df['plot_row'].max() + 1

            fig_height = BASE_HEIGHT + TIMELINE_SPACING * n_rows_this_slide
            fig_width = 10
            fig, ax = plt.subplots(figsize=(fig_width, fig_height))

            # Draw timeline lines
            for i in range(n_rows_this_slide):
                ax.hlines(y=i, xmin=fy_start, xmax=fy_end, color='tab:blue', linewidth=2)

            # Plot circles using total TCV for each Deal Id
            for _, row in slide_plot_df.iterrows():
                x = row['Close Date']
                y = row['plot_row']
                deal_id = row['Deal Id']
                stage_color = get_stage_color(row.get('Stage', 'NA')).lower()
                
                # Use total TCV for this Deal Id
                atr_total = deal_id_total.get(deal_id, 0)
                circle_size = get_circle_size(atr_total, min_atr, max_atr)
                
                ax.plot(x, y, 'o', color=stage_color, markersize=circle_size)
                
                # Format the Deal Id text
                ax.text(x, y + 0.12, deal_id, fontsize=8, ha='center', va='bottom', rotation=20, fontweight='bold', color='black')

            # Draw month dividers
            months = pd.date_range(start=fy_start, end=fy_end, freq='MS')
            for month_start in months:
                ax.vlines(month_start, ymin=-0.5, ymax=n_rows_this_slide-0.5, color='lightgrey', linestyle='--', linewidth=1)

            ax.set_yticks([])
            ax.set_xlim([fy_start, fy_end])
            ax.set_ylim(-0.5, n_rows_this_slide - 0.5)
            ax.set_title("")
            ax.set_xlabel('Close Date')
            ax.xaxis.set_major_locator(MonthLocator())
            ax.xaxis.set_major_formatter(DateFormatter('%b\n%Y'))

            add_timeline_legend(ax)

            plt.tight_layout()
            img_stream = io.BytesIO()
            plt.savefig(img_stream, format='png')
            plt.close(fig)
            img_stream.seek(0)
            image_aspect_ratio = fig_width / fig_height if fig_height > 0 else (16.0 / 9.0)

            # Add to PowerPoint
            timeline_slide = prs.slides.add_slide(prs.slide_layouts[1])
            remove_all_placeholders(timeline_slide)
            pic_left, pic_top, pic_width, pic_height = add_content_title(
                timeline_slide,
                prs,
                f"{account} {title_prefix} Timeline" + (f" (page {slide_idx+1})" if n_timelines > MAX_TIMELINES_PER_SLIDE else ""),
                font_size=15,
            )
            add_fitted_picture(timeline_slide, img_stream, pic_left, pic_top, pic_width, pic_height, image_aspect_ratio)

            # --- Notes Section ---
            notes_lines = []

            # Aggregated Detailed Entries (Unique by Deal Id)
            notes_lines.append("Detailed Entries:")
            # Group the DataFrame by Deal Id and Close Date to ensure uniqueness
            unique_deals = acc_slide_df.drop_duplicates(subset=['Deal Id'])
            for month_str, group in unique_deals.groupby(acc_slide_df['Close Date'].dt.to_period('M')):
                month_label = month_str.strftime('%b %Y')
                notes_lines.append(f"{month_label}:")
                for _, row in group.iterrows():
                    deal_id = row['Deal Id']
                    atr_total = int(deal_id_total.get(deal_id, 0))  # Total aggregated value
                    opp_name = row.get('Opportunity Name', "")
                    opp_owner = row.get('Opportunity Owner', "")
                    stage = row.get('Stage', "")
                    close_date = row['Close Date'].strftime('%Y-%m-%d')
                    notes_lines.append(f"- {deal_id}: {atr_total}K | {opp_name} | {opp_owner} | {stage} | {close_date}")

            notes_text = '\n'.join(notes_lines)
            timeline_slide.notes_slide.notes_text_frame.text = notes_text


def add_monthly_timeline_slides(prs, filtered, fy_start, fy_end):
    months = pd.date_range(start=fy_start, end=fy_end, freq='MS')

    for month_start in months:
        month_end = month_start + pd.offsets.MonthEnd(0)

        # Filter data for the specific month
        month_df = filtered[(filtered['Close Date'] >= month_start) & (filtered['Close Date'] <= month_end)].copy()

        if month_df.empty:
            continue

        # Aggregate by Deal Id for total TCV calculation
        deal_id_total = month_df.groupby('Deal Id')['Expected Amount TCV ($000s)'].sum().to_dict()

        # Determine min and max TCV for circle sizing
        if deal_id_total:
            min_atr = min(deal_id_total.values())
            max_atr = max(deal_id_total.values())
        else:
            min_atr = max_atr = 1.0

        # Prepare slide title
        month_year = month_start.strftime('%B %Y')
        slide_title = f"New Opportunities in {month_year}"

        # Create a new slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        remove_all_placeholders(slide)
        pic_left, pic_top, pic_width, pic_height = add_content_title(slide, prs, slide_title, font_size=16)

        # Prepare the plot
        fig_width = 10
        fig_height = 5
        fig, ax = plt.subplots(figsize=(fig_width, fig_height))
        ax.set_title("")
        ax.set_xlabel('Close Date')
        ax.xaxis.set_major_locator(MonthLocator())
        ax.xaxis.set_major_formatter(DateFormatter('%b\n%Y'))

        # Plot timeline
        for idx, (deal_id, group) in enumerate(month_df.groupby('Deal Id')):
            x = group['Close Date']
            y = [idx] * len(x)
            atr_total = deal_id_total[deal_id]

            # Determine circle size
            circle_size = get_circle_size(atr_total, min_atr, max_atr)

            # Determine stage color
            stage_color = get_stage_color(group['Stage'].iloc[0]).lower()

            ax.plot(x, y, 'o', color=stage_color, markersize=circle_size)
            # Add Deal Id text
            ax.text(x.iloc[0], y[0] + 0.1, deal_id, fontsize=8, ha='center', va='bottom', rotation=20, fontweight='bold', color='black')

        ax.set_yticks([])
        ax.set_xlim([month_start, month_end + pd.Timedelta(days=1)])  # +1 day for aesthetics
        ax.set_ylim(-0.5, idx + 0.5)

        # Add legend at the top
        add_timeline_legend(ax)

        plt.tight_layout(rect=[0, 0, 0.85, 1])  # Leave space for legend
        img_stream = io.BytesIO()
        plt.savefig(img_stream, format='png')
        plt.close(fig)
        img_stream.seek(0)
        image_aspect_ratio = fig_width / fig_height if fig_height > 0 else (16.0 / 9.0)
        add_fitted_picture(slide, img_stream, pic_left, pic_top, pic_width, pic_height, image_aspect_ratio)

        # --- Notes Section ---
        notes_lines = []

        # Aggregated Detailed Entries (Unique by Deal Id)
        notes_lines.append("Detailed Entries:")
        unique_deals = month_df.drop_duplicates(subset=['Deal Id'])
        for _, row in unique_deals.iterrows():
            deal_id = row['Deal Id']
            account_name = row.get('Account Name', "")
            atr_total = int(deal_id_total.get(deal_id, 0))  # Total aggregated value
            opp_name = row.get('Opportunity Name', "")
            opp_owner = row.get('Opportunity Owner', "")
            stage = row.get('Stage', "")
            close_date = row['Close Date'].strftime('%Y-%m-%d')
            notes_lines.append(f"- {deal_id}: {account_name} | {atr_total}K | {opp_name} | {opp_owner} | {stage} | {close_date}")

        notes_text = '\n'.join(notes_lines)
        slide.notes_slide.notes_text_frame.text = notes_text



def create_ppt_for_filter(prs, filtered, fy_start, fy_end, initial_fy, final_fy, file_creation_date, title_prefix, columns, min_tcv=0):
    customer_names = get_customer_names(filtered)
    customer_scope = format_customer_scope(customer_names, max_names=5)
    # --- TITLE SLIDE ---
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = title_slide.shapes.title
    title_text = f"{title_prefix} in {customer_scope}"
    if title_shape is not None:
        title_shape.text = title_text
    else:
        title_box = title_slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8.8), Inches(0.8))
        title_box.text_frame.text = title_text
    subtitle_text = (
        f"Date Range: {initial_fy} – {final_fy}\n"
        f"Fiscal Dates: {fy_start.strftime('%b %d, %Y')} – {fy_end.strftime('%b %d, %Y')}\n"
        f"File Created: {file_creation_date}"
    )

    if customer_names:
        subtitle_text += f"\nCustomers ({len(customer_names)}): {format_customer_scope(customer_names, max_names=8)}"
    else:
        subtitle_text += "\nCustomers: Not provided"
    if min_tcv > 0:
        subtitle_text += f"\nMinimum Expected TCV: ${int(min_tcv):,}"
    try:
        subtitle_shape = title_slide.placeholders[1]
        subtitle_shape.text = subtitle_text
        subtitle_tf = subtitle_shape.text_frame
    except Exception:
        subtitle_box = title_slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.5), Inches(1.8))
        subtitle_tf = subtitle_box.text_frame
        subtitle_tf.text = subtitle_text
    
    # Reduce subtitle font size
    for paragraph in subtitle_tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(14)

    # Add legend text box on title slide
    if min_tcv > 0:
        # Position the legend at the bottom of the slide
        legend_left = Inches(1)
        legend_top = Inches(6)
        legend_width = Inches(8)
        legend_height = Inches(1)
        
        legend_box = title_slide.shapes.add_textbox(legend_left, legend_top, legend_width, legend_height)
        legend_frame = legend_box.text_frame
        legend_frame.word_wrap = True
        
        # Add legend title
        p = legend_frame.paragraphs[0]
        p.text = "Note:"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        
        # Add explanation
        p = legend_frame.add_paragraph()
        p.text = f"This presentation includes only deals with total aggregated Expected TCV >= ${int(min_tcv)}K"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_before = Pt(6)

    # --- SUMMARY TABLE SLIDES ---
    value_col = 'Expected Amount TCV ($000s)'
    filtered_scope = filtered.copy()
    filtered_scope['Customer Name'] = filtered_scope['Customer Name'].fillna('').astype(str).str.strip()
    add_summary_table_slide(
        prs, filtered_scope, fy_start, fy_end, value_col, summary_scope='All Customers'
    )
    for customer in customer_names:
        customer_df = filtered_scope[filtered_scope['Customer Name'] == customer]
        add_summary_table_slide(
            prs, customer_df, fy_start, fy_end, value_col, summary_scope=customer
        )
    # --- Table and Timeline slides ---
    add_table_slides(prs, filtered, columns, title_prefix)
    add_timeline_slides(prs, filtered, fy_start, fy_end, title_prefix, min_tcv=min_tcv)




def create_new_ops_ppt(initial_fy, final_fy, excel_filename, min_tcv=0, template_pptx=None):
    fy_start, fy_end = get_fy_range(initial_fy, final_fy)
    df = check_excel_file(excel_filename)
    if len(df) == 0:
        print("Warning: The input Excel file contains no data.")
    df['Close Date'] = pd.to_datetime(df['Close Date'], errors='coerce')
    filtered = df[(df['Close Date'] >= fy_start) & (df['Close Date'] <= fy_end)].copy()
    if filtered.empty:
        print("Warning: No new opportunities found within the selected date range.")
    filtered = filtered[BASE_SLIDE_COLUMNS + ['Customer Name']]

    if 'Deal Id' in filtered.columns:
        filtered['Deal Id'] = filtered['Deal Id'].apply(
            lambda x: str(int(x)) if pd.notnull(x) and isinstance(x, (float, int)) and float(x) == int(x) else str(x)
        )

    for col in ['Expected Amount TCV ($000s)']:
        filtered[col] = (
            filtered[col]
            .astype(str)
            .str.replace(r'[\$,]', '', regex=True)
            .replace('nan', '0')
            .astype(float)
        )
    
    # Aggregate totals and filter based on min_tcv
    deal_id_total = filtered.groupby('Deal Id')['Expected Amount TCV ($000s)'].sum().to_dict()
    filtered_deal_ids = {deal_id for deal_id, total in deal_id_total.items() if total >= min_tcv}
    filtered = filtered[filtered['Deal Id'].isin(filtered_deal_ids)]

    file_creation_date = get_file_creation_date(excel_filename)

    # Format min_tcv to remove '.0' if it's an integer
    min_tcv_str = f"{min_tcv:.0f}" if min_tcv % 1 == 0 else f"{min_tcv}"

    fy_suffix = f"_{initial_fy}-{final_fy}_TCV_MIN_{min_tcv_str}"

    # --- NEW OPS PPT ---
    prs = Presentation(template_pptx) if template_pptx else Presentation()  # Initialize the presentation once
    filtered_product = filtered[filtered['Expected Amount TCV ($000s)'] > 0]
    product_columns = [c for c in BASE_SLIDE_COLUMNS]
    if not filtered_product.empty:
        create_ppt_for_filter(
            prs, filtered_product, fy_start, fy_end,
            initial_fy, final_fy, file_creation_date,
            "New Opportunities", product_columns, min_tcv=min_tcv
        )
        # Add monthly slides at the end
        add_monthly_timeline_slides(prs, filtered_product, fy_start, fy_end)
    else:
        print("No new opportunities found.")

    # Save the presentation
    pptx_filename = os.path.splitext(excel_filename)[0] + fy_suffix + '.pptx'
    prs.save(pptx_filename)
    print(f"Presentation saved as: {pptx_filename}")

if __name__ == '__main__':
    parser = argparse.ArgumentParser(description="Generate Cisco new opportunities PowerPoint from Excel.")
    parser.add_argument('initial_fy', type=str, help="Initial date (format: QQFYXX, e.g., Q1FY26)")
    parser.add_argument('final_fy', type=str, help="Final date (format: QQFYXX, e.g., Q3FY26)")
    parser.add_argument('excel_filename', type=str, help="Input Excel file name (.xlsx)")
    parser.add_argument('--min-tcv', type=float, default=0, help="Minimum Expected Amount TCV ($000s)")
    parser.add_argument('--template-pptx', type=str, dest='template_pptx', default=None, help="Optional PowerPoint template file (.pptx)")
    args = parser.parse_args()

    # Robust parameter validation
    if not robust_validate_fy_quarter(args.initial_fy):
        sys.exit(2)
    if not robust_validate_fy_quarter(args.final_fy):
        sys.exit(2)
    if not robust_check_excel_file(args.excel_filename):
        sys.exit(2)
    if not robust_check_template_file(args.template_pptx):
        sys.exit(2)

    try:
        create_new_ops_ppt(args.initial_fy, args.final_fy, args.excel_filename, args.min_tcv, args.template_pptx)
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(2)

